#!/usr/bin/env python

import os
import sys
import logging
from datetime import datetime
from typing import Dict, List, Any, Optional
from pathlib import Path
import json
import traceback
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# Add project root to path
project_root = os.path.abspath(os.path.join(os.path.dirname(__file__), '../..'))
sys.path.append(project_root)

# Import existing components
from Sample_Data.vector_store.storage_manager import StorageManager
from crypto_qa import EnhancedCryptoQA
from agentic_rag import CryptoDueDiligenceSystem
from .utils.chart_generator import ChartGenerator
from .utils.style_manager import StyleManager
from .utils.data_formatter import DataFormatter
from .report_config import REPORT_TOPICS, TEMPLATE_SETTINGS

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

class CryptoPPTGenerator:
    """
    Generates PowerPoint presentations from crypto due diligence data
    using existing Weaviate collections and Q&A system.
    """
    
    def __init__(self, template_path: Optional[str] = None):
        self.storage = StorageManager()
        self.qa_system = EnhancedCryptoQA()
        self.due_diligence_system = CryptoDueDiligenceSystem()
        self.chart_generator = ChartGenerator()
        self.style_manager = StyleManager()
        self.data_formatter = DataFormatter()
        
        # Default template or provided path
        if template_path and os.path.exists(template_path):
            self.template_path = template_path
        else:
            self.template_path = None
        
        # Create reports directory
        self.reports_dir = os.path.join(os.path.dirname(__file__), 'reports')
        os.makedirs(self.reports_dir, exist_ok=True)
        logger.info(f"Reports directory created/verified at: {self.reports_dir}")
        
        # Initialize components
        self._initialize_components()
    
    def _initialize_components(self):
        """Initialize necessary components."""
        try:
            if not self.storage.connect():
                logger.error("Failed to connect to Weaviate storage")
            
            if not self.due_diligence_system.initialize():
                logger.error("Failed to initialize Due Diligence System")
                
            logger.info("Report generator initialized successfully")
        except Exception as e:
            logger.error(f"Error initializing components: {str(e)}")
            logger.error(traceback.format_exc())
    
    def generate_report(self, topics: List[str], crypto_symbol: str, output_filename: Optional[str] = None) -> str:
        """
        Generate a PowerPoint report for a specific cryptocurrency.
        
        Args:
            topics: List of due diligence topics to cover
            crypto_symbol: Cryptocurrency symbol (e.g., 'BTC', 'ETH')
            output_filename: Optional filename for the report
            
        Returns:
            str: Path to the generated report file
        """
        try:
            # Create presentation
            if self.template_path:
                prs = Presentation(self.template_path)
            else:
                prs = Presentation()
            
            # Apply styling
            self.style_manager.apply_theme(prs)
            
            # Add slides
            self._add_title_slide(prs, crypto_symbol)
            self._add_executive_summary_slide(prs, crypto_symbol)
            
            # Generate slides for each topic
            for topic in topics:
                if topic in REPORT_TOPICS:
                    self._generate_topic_slide(prs, topic, crypto_symbol)
            
            # Add data slides
            self._add_market_data_slide(prs, crypto_symbol)
            self._add_price_trend_slide(prs, crypto_symbol)
            self._add_sentiment_slide(prs, crypto_symbol)
            self._add_risk_assessment_slide(prs, crypto_symbol)
            
            # Add conclusion
            self._add_summary_slide(prs, crypto_symbol)
            
            # Generate output filename if not provided
            if not output_filename:
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                output_filename = f"{crypto_symbol}_due_diligence_report_{timestamp}.pptx"
            
            # Ensure the filename ends with .pptx
            if not output_filename.endswith('.pptx'):
                output_filename += '.pptx'
            
            # Create full output path in reports directory
            output_path = os.path.join(self.reports_dir, output_filename)
            
            # Save presentation
            prs.save(output_path)
            logger.info(f"Report saved to {output_path}")
            
            return output_path
            
        except Exception as e:
            logger.error(f"Error generating report: {str(e)}")
            logger.error(traceback.format_exc())
            return None
    
    def _add_title_slide(self, prs, crypto_symbol):
        """Add a title slide to the presentation."""
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        # Set title text
        title.text = f"{crypto_symbol} Due Diligence Report"
        self.style_manager.apply_title_formatting(title)
        
        # Set subtitle text
        subtitle.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}\nBased on Multi-Source Data Analysis"
        self.style_manager.apply_subtitle_formatting(subtitle)
    
    def _add_executive_summary_slide(self, prs, crypto_symbol):
        """Add an executive summary slide."""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Executive Summary"
        self.style_manager.apply_title_formatting(title)
        
        # Generate executive summary using Q&A system
        question = f"Provide a brief executive summary for {crypto_symbol} including key highlights about market position, technology, and investment potential."
        summary = self.qa_system.answer_question(question)
        
        # Add summary text
        content = slide.placeholders[1]
        content.text = summary
        self.style_manager.apply_body_formatting(content)
    
    def _generate_topic_slide(self, prs, topic, crypto_symbol):
        """Generate a slide for a specific due diligence topic."""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Format topic for display
        formatted_topic = topic.replace('_', ' ').title()
        
        title = slide.shapes.title
        title.text = f"{formatted_topic} Analysis"
        self.style_manager.apply_title_formatting(title)
        
        # Get topic-specific questions
        if topic == "legal_regulatory":
            question = f"What is the current legal and regulatory status of {crypto_symbol}? Include information about compliance, licenses, and regulatory challenges."
        elif topic == "team_background":
            question = f"Who are the key team members behind {crypto_symbol}? Include information about their background, experience, and track record."
        elif topic == "technical":
            question = f"What are the technical aspects of {crypto_symbol}? Discuss its blockchain architecture, consensus mechanism, and technical innovations."
        elif topic == "financial":
            question = f"What is the financial structure of {crypto_symbol}? Include tokenomics, funding, revenue model, and financial metrics."
        elif topic == "governance":
            question = f"How is {crypto_symbol} governed? Discuss the decision-making process, voting mechanisms, and community participation."
        elif topic == "risk":
            question = f"What are the key risks associated with {crypto_symbol}? Include technical, regulatory, market, and operational risks."
        elif topic == "on_chain":
            question = f"What do on-chain analytics reveal about {crypto_symbol}? Include network activity, whale movements, and on-chain metrics."
        else:
            question = f"Provide a comprehensive analysis of {topic} for {crypto_symbol}."
        
        # Get answer from Q&A system
        answer = self.qa_system.answer_question(question)
        
        # Add content
        content = slide.placeholders[1]
        content.text = answer
        self.style_manager.apply_body_formatting(content)
    
    def _add_market_data_slide(self, prs, crypto_symbol):
        """Add a slide with market data."""
        slide_layout = prs.slide_layouts[5]  # Title only layout
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = f"{crypto_symbol} Market Data"
        self.style_manager.apply_title_formatting(title)
        
        # Retrieve market data
        symbol_formatted = f"{crypto_symbol}USDT"
        market_data = self.storage.retrieve_market_data(symbol_formatted, limit=1)
        
        if market_data and len(market_data) > 0:
            data = market_data[0]
            
            # Create market data table
            self._create_market_data_table(slide, data)
            
            # Add data explanation
            textbox = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
            textbox.text = f"Market data as of {data.get('timestamp', 'N/A')}"
            self.style_manager.apply_body_formatting(textbox)
    
    def _add_price_trend_slide(self, prs, crypto_symbol):
        """Add a slide with price trend chart."""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = f"{crypto_symbol} Price Trend"
        self.style_manager.apply_title_formatting(title)
        
        # Retrieve historical data
        symbol_formatted = f"{crypto_symbol}USDT"
        historical_data = self.storage.retrieve_market_data(symbol_formatted, limit=30)
        
        if historical_data and len(historical_data) > 0:
            # Generate price chart
            chart_stream = self.chart_generator.generate_price_chart(historical_data)
            
            # Add chart to slide
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(4)
            pic = slide.shapes.add_picture(chart_stream, left, top, width=width, height=height)
    
    def _add_sentiment_slide(self, prs, crypto_symbol):
        """Add a slide with sentiment analysis."""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = f"{crypto_symbol} Sentiment Analysis"
        self.style_manager.apply_title_formatting(title)
        
        # Get sentiment data
        sentiment_data = self.storage.get_sentiment_stats(crypto_symbol.lower())
        
        if sentiment_data and "sentiment_distribution" in sentiment_data:
            # Generate sentiment pie chart
            chart_stream = self.chart_generator.generate_sentiment_pie_chart(
                sentiment_data["sentiment_distribution"]
            )
            
            # Add chart to slide
            left = Inches(1)
            top = Inches(2)
            width = Inches(6)
            height = Inches(4)
            pic = slide.shapes.add_picture(chart_stream, left, top, width=width, height=height)
            
            # Add sentiment insights
            textbox = slide.shapes.add_textbox(Inches(7), Inches(2), Inches(2.5), Inches(4))
            text_frame = textbox.text_frame
            text_frame.word_wrap = True
            
            p = text_frame.add_paragraph()
            p.text = f"Total Articles: {sentiment_data.get('total_articles', 0)}"
            
            p = text_frame.add_paragraph()
            p.text = f"Average Sentiment: {sentiment_data.get('avg_sentiment', 0):.2f}"
            
            p = text_frame.add_paragraph()
            p.text = f"Trend: {sentiment_data.get('trend', 'unknown')}"
            
            self.style_manager.apply_body_formatting(textbox)
    
    def _add_risk_assessment_slide(self, prs, crypto_symbol):
        """Add a slide with risk assessment."""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = f"{crypto_symbol} Risk Assessment"
        self.style_manager.apply_title_formatting(title)
        
        # Get risk assessment using Q&A system
        question = f"Provide a comprehensive risk assessment for {crypto_symbol} including technical risks, regulatory risks, market risks, and operational risks. Rate each risk as High/Medium/Low and provide mitigation strategies."
        risk_answer = self.qa_system.answer_question(question)
        
        content = slide.placeholders[1]
        content.text = risk_answer
        self.style_manager.apply_body_formatting(content)
    
    def _add_summary_slide(self, prs, crypto_symbol):
        """Add a summary slide."""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = f"{crypto_symbol} Investment Conclusion"
        self.style_manager.apply_title_formatting(title)
        
        # Generate summary using Q&A system
        question = f"Based on the due diligence analysis, provide an investment conclusion for {crypto_symbol}. Include key strengths, main concerns, and overall recommendation."
        summary = self.qa_system.answer_question(question)
        
        content = slide.placeholders[1]
        content.text = summary
        self.style_manager.apply_body_formatting(content)
    
    def _create_market_data_table(self, slide, data):
        """Create a table for market data."""
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(3)
        
        # Create table
        shape = slide.shapes.add_table(6, 2, left, top, width, height)
        table = shape.table
        
        # Style table
        self.style_manager.apply_table_formatting(table)
        
        # Define metrics and their values
        metrics = [
            ("Metric", "Value"),
            ("Price", f"${data.get('price', 'N/A'):,.2f}"),
            ("24h Change", f"{data.get('price_change_24h', 'N/A')}%"),
            ("Market Cap", f"${data.get('market_cap', 'N/A'):,.0f}"),
            ("24h Volume", f"${data.get('volume_24h', 'N/A'):,.0f}"),
            ("Timestamp", data.get('timestamp', 'N/A'))
        ]
        
        # Fill table
        for row_idx, (metric, value) in enumerate(metrics):
            table.cell(row_idx, 0).text = metric
            table.cell(row_idx, 1).text = value
            
            # Format header row
            if row_idx == 0:
                for col_idx in range(2):
                    cell = table.cell(row_idx, col_idx)
                    self.style_manager.apply_table_header_formatting(cell)
    
    def close(self):
        """Close connections."""
        if hasattr(self, 'storage') and self.storage:
            self.storage.close()
        if hasattr(self, 'qa_system') and self.qa_system:
            self.qa_system.close()