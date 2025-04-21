"""
Style management for PowerPoint reports
"""
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.dml import MSO_THEME_COLOR
from typing import Dict
import logging

logger = logging.getLogger(__name__)

class StyleManager:
    """Manage styles for PowerPoint presentations."""
    
    def __init__(self):
        # Load configuration from report_config
        try:
            from ..report_config import TEMPLATE_SETTINGS, TABLE_SETTINGS
            self.template_settings = TEMPLATE_SETTINGS
            self.table_settings = TABLE_SETTINGS
        except ImportError:
            logger.warning("Could not import settings, using defaults")
            self.template_settings = {
                "font_family": "Arial",
                "title_font_size": Pt(44),
                "subtitle_font_size": Pt(32),
                "body_font_size": Pt(18),
                "caption_font_size": Pt(14),
                "primary_color": RGBColor(0, 84, 164),
                "secondary_color": RGBColor(127, 127, 127),
                "text_color": RGBColor(0, 0, 0)
            }
            self.table_settings = {
                "header_fill": RGBColor(0, 84, 164),
                "header_text_color": RGBColor(255, 255, 255),
                "row_fill_1": RGBColor(255, 255, 255),
                "row_fill_2": RGBColor(242, 242, 242),
                "border_color": RGBColor(127, 127, 127),
                "border_width": Pt(0.5)
            }
    
    def apply_theme(self, presentation):
        """Apply the theme to the entire presentation."""
        try:
            # Apply font family to the entire presentation
            slide_master = presentation.slide_master
            
            for shape in slide_master.shapes:
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.font.name = self.template_settings["font_family"]
            
            logger.info("Theme applied successfully")
        except Exception as e:
            logger.error(f"Error applying theme: {str(e)}")
    
    def apply_title_formatting(self, title_shape):
        """Apply title formatting to a shape."""
        try:
            if hasattr(title_shape, "text_frame"):
                text_frame = title_shape.text_frame
                
                for paragraph in text_frame.paragraphs:
                    paragraph.font.name = self.template_settings["font_family"]
                    paragraph.font.size = self.template_settings["title_font_size"]
                    paragraph.font.color.rgb = self.template_settings["primary_color"]
                    paragraph.font.bold = True
                    paragraph.alignment = PP_ALIGN.CENTER
                
                # Center the text vertically
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception as e:
            logger.error(f"Error applying title formatting: {str(e)}")
    
    def apply_subtitle_formatting(self, subtitle_shape):
        """Apply subtitle formatting to a shape."""
        try:
            if hasattr(subtitle_shape, "text_frame"):
                text_frame = subtitle_shape.text_frame
                
                for paragraph in text_frame.paragraphs:
                    paragraph.font.name = self.template_settings["font_family"]
                    paragraph.font.size = self.template_settings["subtitle_font_size"]
                    paragraph.font.color.rgb = self.template_settings["secondary_color"]
                    paragraph.alignment = PP_ALIGN.CENTER
                
                # Center the text vertically
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception as e:
            logger.error(f"Error applying subtitle formatting: {str(e)}")
    
    def apply_body_formatting(self, body_shape):
        """Apply body text formatting to a shape."""
        try:
            if hasattr(body_shape, "text_frame"):
                text_frame = body_shape.text_frame
                text_frame.word_wrap = True
                
                for paragraph in text_frame.paragraphs:
                    paragraph.font.name = self.template_settings["font_family"]
                    paragraph.font.size = self.template_settings["body_font_size"]
                    paragraph.font.color.rgb = self.template_settings["text_color"]
                    paragraph.alignment = PP_ALIGN.LEFT
                    
                    # Adjust line spacing
                    paragraph.line_spacing = 1.2
                
                # Set text box to resize shape to fit text
                text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        except Exception as e:
            logger.error(f"Error applying body formatting: {str(e)}")
    
    def apply_table_formatting(self, table):
        """Apply formatting to a table."""
        try:
            # Format the entire table
            for row_idx, row in enumerate(table.rows):
                for col_idx, cell in enumerate(row.cells):
                    # Set cell margins
                    cell.margin_left = Inches(0.1)
                    cell.margin_right = Inches(0.1)
                    cell.margin_top = Inches(0.05)
                    cell.margin_bottom = Inches(0.05)
                    
                    # Apply alternating row colors
                    if row_idx > 0:  # Skip header row
                        if row_idx % 2 == 0:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = self.table_settings["row_fill_2"]
                        else:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = self.table_settings["row_fill_1"]
                    
                    # Format text in cell
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.name = self.template_settings["font_family"]
                        paragraph.font.size = self.template_settings["body_font_size"]
                        paragraph.font.color.rgb = self.template_settings["text_color"]
                        
                        # Align text based on column (first column left, others right)
                        if col_idx == 0:
                            paragraph.alignment = PP_ALIGN.LEFT
                        else:
                            paragraph.alignment = PP_ALIGN.RIGHT
            
            # Set table borders
            self._set_table_borders(table)
            
        except Exception as e:
            logger.error(f"Error applying table formatting: {str(e)}")
    
    def apply_table_header_formatting(self, cell):
        """Apply header formatting to a table cell."""
        try:
            # Set header background color
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.table_settings["header_fill"]
            
            # Format header text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = self.template_settings["font_family"]
                paragraph.font.size = self.template_settings["body_font_size"]
                paragraph.font.color.rgb = self.table_settings["header_text_color"]
                paragraph.font.bold = True
                paragraph.alignment = PP_ALIGN.CENTER
        except Exception as e:
            logger.error(f"Error applying table header formatting: {str(e)}")
    
    def _set_table_borders(self, table):
        """Set borders for a table."""
        try:
            # This is a simplified version - PowerPoint table border API is complex
            # For a production version, you would need more detailed border control
            for row in table.rows:
                for cell in row.cells:
                    # Set borders (requires python-pptx extensions for full control)
                    pass
        except Exception as e:
            logger.error(f"Error setting table borders: {str(e)}")
    
    def apply_chart_style(self, chart):
        """Apply style to a chart."""
        try:
            # Set chart style
            chart.chart_style = 2  # Simple chart style
            
            # Format chart title
            if chart.has_title:
                chart.chart_title.text_frame.paragraphs[0].font.name = self.template_settings["font_family"]
                chart.chart_title.text_frame.paragraphs[0].font.size = self.template_settings["body_font_size"]
                chart.chart_title.text_frame.paragraphs[0].font.bold = True
            
            # Format axis titles
            if chart.has_axis:
                if chart.category_axis.has_title:
                    chart.category_axis.axis_title.text_frame.paragraphs[0].font.name = self.template_settings["font_family"]
                    chart.category_axis.axis_title.text_frame.paragraphs[0].font.size = self.template_settings["caption_font_size"]
                
                if chart.value_axis.has_title:
                    chart.value_axis.axis_title.text_frame.paragraphs[0].font.name = self.template_settings["font_family"]
                    chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = self.template_settings["caption_font_size"]
            
        except Exception as e:
            logger.error(f"Error applying chart style: {str(e)}")