from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime
import os
import re

def sanitize_filename(name):
    return re.sub(r'[\\/*?:"<>|]', "", name).replace(" ", "_")

def create_crypto_ppt_styled(topic: str, qa_summary: str, weaviate_insights: str, output_dir: str = "reports") -> str:
    prs = Presentation()
    
    # Set default font
    for slide_layout in prs.slide_layouts:
        for placeholder in slide_layout.placeholders:
            if hasattr(placeholder, "text_frame") and placeholder.text_frame:
                for paragraph in placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Calibri'
                        run.font.size = Pt(18)
                        run.font.bold = True

    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide_layout)
    slide1.shapes.title.text = f"Weaviate Insights: {topic}"
    slide1.placeholders[1].text = "Automated Due Diligence Report"

    # Insights Slide
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Contextual Insights"
    text_frame = slide2.placeholders[1].text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = weaviate_insights.strip()
    p.font.size = Pt(14)

    # Q&A Slide
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Q&A Summary"
    text_frame = slide3.placeholders[1].text_frame
    text_frame.clear()
    qas = qa_summary.strip().split('\n')
    for qa in qas:
        para = text_frame.add_paragraph()
        para.text = qa.strip()
        para.level = 0
        para.font.size = Pt(14)

    # Save file
    os.makedirs(output_dir, exist_ok=True)
    filename = f"{sanitize_filename(topic)}_styled_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    filepath = os.path.join(output_dir, filename)
    prs.save(filepath)
    return filepath
