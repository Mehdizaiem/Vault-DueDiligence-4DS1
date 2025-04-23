"""
Chart Factory Module

This module handles the creation of various chart types for PowerPoint presentations.
It provides a clean interface for generating charts with consistent styling and
formatting based on due diligence data.
"""

import os
import logging
from typing import Dict, List, Any, Optional, Union, Tuple
import math
import random

from pptx.chart.data import CategoryChartData, ChartData, BubbleChartData
from pptx.chart.data import XyChartData, XySeriesData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from reporting.design_elements import DesignElements

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

class ChartFactory:
    """
    Creates professional data visualizations for PowerPoint presentations
    with consistent styling, appropriate chart types, and financial formatting.
    """
    
    def __init__(self,):
        """Initialize the chart factory with design elements."""
        self.design = DesignElements()
        
    def create_pie_chart(self, slide, chart_data: List[Tuple[str, float]], 
                         left: float, top: float, width: float, height: float,
                         title: Optional[str] = None, 
                         has_labels: bool = True) -> None:
        """
        Create a pie chart for portfolio allocation or other percentage-based data.
        
        Args:
            slide: Slide to add the chart to
            chart_data: List of (label, value) tuples
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            title: Chart title (optional)
            has_labels: Whether to show data labels
        """
        # Create chart data
        chart_data_obj = CategoryChartData()
        chart_data_obj.categories = [item[0] for item in chart_data]
        
        # Add series
        series_values = [item[1] for item in chart_data]
        chart_data_obj.add_series('Allocation', series_values)
        
        # Add chart to slide
        x, y = Inches(left), Inches(top)
        cx, cy = Inches(width), Inches(height)
        
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data_obj
        ).chart
        
        # Set title if provided
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
        else:
            chart.has_title = False
        
        # Configure legend
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
        
        # Configure data labels
        chart.plots[0].has_data_labels = has_labels
        data_labels = chart.plots[0].data_labels
        data_labels.font.size = Pt(10)
        data_labels.font.color.rgb = RGBColor(255, 255, 255)
        data_labels.position = XL_LABEL_POSITION.CENTER
        
        # Show percentage
        data_labels.number_format = '0.0%'
        data_labels.show_percentage = True
        data_labels.show_category_name = False
        
        # Apply custom colors
        series = chart.series[0]
        point_count = len(chart_data)
        
        for idx, point in enumerate(series.points):
            # Use our design color scheme
            if point_count <= len(self.design.CHART_COLORS):
                color = self.design.CHART_COLORS[idx % len(self.design.CHART_COLORS)]
            else:
                # Generate colors if we have more data points than predefined colors
                color_idx = idx % len(self.design.CHART_COLORS)
                base_color = self.design.CHART_COLORS[color_idx]
                # Slight variation of the base color
                r, g, b = base_color
                variation = 30  # Color variation amount
                r = max(0, min(255, r + random.randint(-variation, variation)))
                g = max(0, min(255, g + random.randint(-variation, variation)))
                b = max(0, min(255, b + random.randint(-variation, variation)))
                color = (r, g, b)
                
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = RGBColor(*color)
    
    def create_bar_chart(self, slide, chart_data: List[Tuple[str, float]], 
                        left: float, top: float, width: float, height: float,
                        title: Optional[str] = None, 
                        vertical: bool = True) -> None:
        """
        Create a bar chart for comparing values across categories.
        
        Args:
            slide: Slide to add the chart to
            chart_data: List of (label, value) tuples
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            title: Chart title (optional)
            vertical: Whether bars should be vertical (True) or horizontal (False)
        """
        # Create chart data
        chart_data_obj = CategoryChartData()
        chart_data_obj.categories = [item[0] for item in chart_data]
        
        # Add series
        series_values = [item[1] for item in chart_data]
        chart_data_obj.add_series('Value', series_values)
        
        # Add chart to slide
        x, y = Inches(left), Inches(top)
        cx, cy = Inches(width), Inches(height)
        
        # Choose chart type based on orientation
        chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED if vertical else XL_CHART_TYPE.BAR_CLUSTERED
        
        chart = slide.shapes.add_chart(
            chart_type, x, y, cx, cy, chart_data_obj
        ).chart
        
        # Set title if provided
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
        else:
            chart.has_title = False
        
        # Configure legend
        chart.has_legend = False
        
        # Configure data labels
        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.font.size = Pt(10)
        
        # Apply custom colors
        series = chart.series[0]
        fill_color = self.design.CHART_COLORS[0]
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = RGBColor(*fill_color)
        
        # Set category axis text size
        category_axis = chart.category_axis
        category_axis.tick_labels.font.size = Pt(9)
        
        # Set value axis text size
        value_axis = chart.value_axis
        value_axis.tick_labels.font.size = Pt(9)
    
    def create_radar_chart(self, slide, labels: List[str], values: List[float],
                          left: float, top: float, width: float, height: float,
                          title: Optional[str] = None) -> None:
        """
        Create a radar chart for risk assessment visualization.
        
        Args:
            slide: Slide to add the chart to
            labels: Category labels for the radar points
            values: Values for each category (0-10 scale)
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            title: Chart title (optional)
        """
        # Create chart data
        chart_data_obj = CategoryChartData()
        chart_data_obj.categories = labels
        
        # Add series
        chart_data_obj.add_series('Risk Score', values)
        
        # Add chart to slide
        x, y = Inches(left), Inches(top)
        cx, cy = Inches(width), Inches(height)
        
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.RADAR, x, y, cx, cy, chart_data_obj
        ).chart
        
        # Set title if provided
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
        else:
            chart.has_title = False
        
        # Configure legend
        chart.has_legend = False
        
        # Style the series
        series = chart.series[0]
        
        # Apply line color and style
        series.format.line.solid()
        series.format.line.color.rgb = RGBColor(*self.design.PRIMARY_COLOR)
        series.format.line.width = Pt(2.5)
        
        # Add markers
        series.marker.style = 2  # Circular marker
        series.marker.size = 8
        series.marker.format.fill.solid()
        series.marker.format.fill.fore_color.rgb = RGBColor(*self.design.PRIMARY_COLOR)
        
        # Add data labels
        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.font.size = Pt(9)
        data_labels.font.bold = True
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    
    def create_gauge_chart(self, slide, value: float, 
                          left: float, top: float, width: float, height: float,
                          title: Optional[str] = None,
                          min_value: float = 0.0, max_value: float = 100.0,
                          green_threshold: float = 70.0, 
                          yellow_threshold: float = 30.0) -> None:
        """
        Create a gauge chart for displaying a single metric like risk or compliance score.
        This creates a pie chart styled to look like a gauge.
        
        Args:
            slide: Slide to add the chart to
            value: The value to display (0-100)
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            title: Chart title (optional)
            min_value: Minimum value on the gauge
            max_value: Maximum value on the gauge
            green_threshold: Threshold above which values are considered "good" (green)
            yellow_threshold: Threshold below which values are considered "poor" (red)
        """
        # Normalize value to 0-100 range
        range_size = max_value - min_value
        normalized_value = ((value - min_value) / range_size) * 100 if range_size > 0 else 50
        normalized_value = max(0, min(100, normalized_value))
        
        # Create a pie chart with 2 slices: the value and the remainder
        value_slice = normalized_value
        remainder_slice = 100 - value_slice
        
        # Add "empty" section to make a half-circle
        empty_section = 50  # This makes the gauge a semi-circle
        
        # Create chart data
        chart_data_obj = CategoryChartData()
        chart_data_obj.categories = ['Value', 'Remainder', 'Empty']
        chart_data_obj.add_series('', [value_slice, remainder_slice, empty_section])
        
        # Add chart to slide
        x, y = Inches(left), Inches(top)
        cx, cy = Inches(width), Inches(height)
        
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data_obj
        ).chart
        
        # Set title if provided
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
        else:
            chart.has_title = False
        
        # Hide legend
        chart.has_legend = False
        
        # Configure data labels (show only for the value)
        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        
        # Hide all labels initially
        data_labels.show_percentage = False
        data_labels.show_category_name = False
        data_labels.show_value = False
        
        # Format individual points
        series = chart.series[0]
        
        # Determine color based on thresholds
        if normalized_value >= green_threshold:
            color = self.design.RISK_LOW
        elif normalized_value >= yellow_threshold:
            color = self.design.RISK_MEDIUM
        else:
            color = self.design.RISK_HIGH
        
        # Point 0: Value slice with color based on thresholds
        value_point = series.points[0]
        value_point.format.fill.solid()
        value_point.format.fill.fore_color.rgb = RGBColor(*color)
        
        # Custom label for value
        value_point.data_label.text_frame.text = f"{value:.1f}"
        value_point.data_label.font.size = Pt(24)
        value_point.data_label.font.bold = True
        value_point.data_label.position = XL_LABEL_POSITION.CENTER
        
        # Point 1: Remainder slice (grey)
        remainder_point = series.points[1]
        remainder_point.format.fill.solid()
        remainder_point.format.fill.fore_color.rgb = RGBColor(200, 200, 200)
        
        # Point 2: Empty section (transparent)
        empty_point = series.points[2]
        empty_point.format.fill.background()  # Make it transparent
    
    def create_line_chart(self, slide, 
                         categories: List[str], 
                         series_data: List[Tuple[str, List[float]]],
                         left: float, top: float, width: float, height: float,
                         title: Optional[str] = None,
                         y_axis_title: Optional[str] = None) -> None:
        """
        Create a line chart for time series or trend data with multiple series.
        
        Args:
            slide: Slide to add the chart to
            categories: X-axis categories (e.g., dates, time periods)
            series_data: List of (series_name, values) tuples
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            title: Chart title (optional)
            y_axis_title: Y-axis title (optional)
        """
        # Create chart data
        chart_data_obj = CategoryChartData()
        chart_data_obj.categories = categories
        
        # Add each series
        for series_name, values in series_data:
            chart_data_obj.add_series(series_name, values)
        
        # Add chart to slide
        x, y = Inches(left), Inches(top)
        cx, cy = Inches(width), Inches(height)
        
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE_MARKERS, x, y, cx, cy, chart_data_obj
        ).chart
        
        # Set title if provided
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
        else:
            chart.has_title = False
        
        # Set y-axis title if provided
        if y_axis_title:
            value_axis = chart.value_axis
            value_axis.has_title = True
            value_axis.axis_title.text_frame.text = y_axis_title
        
        # Configure legend
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        
        # Apply custom colors and line styles to each series
        for idx, series in enumerate(chart.series):
            color_idx = idx % len(self.design.CHART_COLORS)
            color = self.design.CHART_COLORS[color_idx]
            
            # Line style
            series.format.line.solid()
            series.format.line.color.rgb = RGBColor(*color)
            series.format.line.width = Pt(2.5)
            
            # Marker style
            series.marker.style = 2  # Circular marker
            series.marker.size = 8
            series.marker.format.fill.solid()
            series.marker.format.fill.fore_color.rgb = RGBColor(*color)
            
            # No data labels by default (can get cluttered with multiple series)
            series.has_data_labels = False
        
        # Adjust axis text size
        category_axis = chart.category_axis
        category_axis.tick_labels.font.size = Pt(9)
        
        value_axis = chart.value_axis
        value_axis.tick_labels.font.size = Pt(9)
    
    def create_bubble_chart(self, slide, 
                          data_points: List[Tuple[str, float, float, float]],
                          left: float, top: float, width: float, height: float,
                          title: Optional[str] = None,
                          x_axis_title: Optional[str] = None,
                          y_axis_title: Optional[str] = None) -> None:
        """
        Create a bubble chart for multi-dimensional data visualization.
        
        Args:
            slide: Slide to add the chart to
            data_points: List of (label, x_value, y_value, size) tuples
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            title: Chart title (optional)
            x_axis_title: X-axis title (optional)
            y_axis_title: Y-axis title (optional)
        """
        # Create chart data
        bubble_chart_data = BubbleChartData()
        
        # Group data points by series (assume one series for simplicity)
        series_name = "Data"
        series = bubble_chart_data.add_series(series_name)
        
        # Add data points
        for label, x, y, size in data_points:
            # Size should be scaled appropriately
            bubble_size = size * 5  # Scale factor for bubble size
            series.add_data_point(x, y, bubble_size)
        
        # Add chart to slide
        x, y = Inches(left), Inches(top)
        cx, cy = Inches(width), Inches(height)
        
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.BUBBLE, x, y, cx, cy, bubble_chart_data
        ).chart
        
        # Set title if provided
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
        else:
            chart.has_title = False
        
        # Set axis titles if provided
        if x_axis_title:
            category_axis = chart.category_axis
            category_axis.has_title = True
            category_axis.axis_title.text_frame.text = x_axis_title
            
        if y_axis_title:
            value_axis = chart.value_axis
            value_axis.has_title = True
            value_axis.axis_title.text_frame.text = y_axis_title
        
        # Configure legend
        chart.has_legend = False
        
        # Apply custom styling to series
        series = chart.series[0]
        color = self.design.CHART_COLORS[0]
        
        # Format markers
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = RGBColor(*color)
        
        # Add semi-transparency
        series.format.fill.transparency = 0.4
        
        # Add data labels with the label text
        series.has_data_labels = True
        
        # This part is tricky as we can't directly set custom text for bubble chart labels
        # In a real implementation, you might need to adjust this part based on PowerPoint's capabilities
        
        # Adjust axis text size
        chart.category_axis.tick_labels.font.size = Pt(9)
        chart.value_axis.tick_labels.font.size = Pt(9)
    
    def create_stacked_bar_chart(self, slide,
                               categories: List[str],
                               series_data: List[Tuple[str, List[float]]],
                               left: float, top: float, width: float, height: float,
                               title: Optional[str] = None,
                               vertical: bool = True) -> None:
        """
        Create a stacked bar chart for showing composition across categories.
        
        Args:
            slide: Slide to add the chart to
            categories: Category labels
            series_data: List of (series_name, values) tuples 
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            title: Chart title (optional)
            vertical: Whether bars should be vertical (True) or horizontal (False)
        """
        # Create chart data
        chart_data_obj = CategoryChartData()
        chart_data_obj.categories = categories
        
        # Add each series
        for series_name, values in series_data:
            chart_data_obj.add_series(series_name, values)
        
        # Add chart to slide
        x, y = Inches(left), Inches(top)
        cx, cy = Inches(width), Inches(height)
        
        # Choose chart type based on orientation
        chart_type = XL_CHART_TYPE.COLUMN_STACKED if vertical else XL_CHART_TYPE.BAR_STACKED
        
        chart = slide.shapes.add_chart(
            chart_type, x, y, cx, cy, chart_data_obj
        ).chart
        
        # Set title if provided
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
        else:
            chart.has_title = False
        
        # Configure legend
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        
        # Apply custom colors to each series
        for idx, series in enumerate(chart.series):
            color_idx = idx % len(self.design.CHART_COLORS)
            color = self.design.CHART_COLORS[color_idx]
            
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = RGBColor(*color)
        
        # Adjust axis text size
        chart.category_axis.tick_labels.font.size = Pt(9)
        chart.value_axis.tick_labels.font.size = Pt(9)