# Path: reporting/pptx_builder.py
import os
import logging
import math
from typing import Dict, List, Any, Optional, Union, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.text.text import _Paragraph, TextFrame
from pptx.shapes.autoshape import Shape
from pptx.table import _Cell

# Assuming these exist in the specified locations
from reporting.chart_factory import ChartFactory
from reporting.design_elements import DesignElements

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

class PresentationBuilder:
    """
    Enhanced builder for creating sophisticated PowerPoint presentations
    with consistent styling, robust error handling, and integration points for LLM narratives.
    """

    def __init__(self, template_path: Optional[str] = None):
        self.design = DesignElements()
        self.chart_factory = ChartFactory()

        # Determine the template path
        if not template_path:
            # If no template path provided, use the default one
            project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
            default_template_path = os.path.join(project_root, "reporting", "templates", "base_template.pptx")
            template_path = default_template_path
            logger.info(f"Using default template: {template_path}")

        try:
            if template_path and os.path.exists(template_path):
                self.prs = Presentation(template_path)
                logger.info(f"Using template: {template_path}")
            else:
                # If template doesn't exist, create a new presentation
                logger.warning(f"Template not found at {template_path}. Using default presentation template.")
                self.prs = Presentation()
                # Set default slide size (Widescreen 16:9)
                self.prs.slide_width = Emu(12192000)  # Inches(13.33)
                self.prs.slide_height = Emu(6858000)  # Inches(7.5)
        except Exception as e:
            logger.error(f"Error initializing Presentation object: {e}", exc_info=True)
            logger.warning("Falling back to a new default presentation.")
            self.prs = Presentation()
            self.prs.slide_width = Emu(12192000)
            self.prs.slide_height = Emu(6858000)

        self.slide_width_emu = self.prs.slide_width
        self.slide_height_emu = self.prs.slide_height
        self.slide_count = len(self.prs.slides)
        logger.info(f"Initialized PresentationBuilder with {self.slide_count} existing slides.")

    def save(self, output_path: str) -> None:
        try:
            output_dir = os.path.dirname(os.path.abspath(output_path))
            if output_dir:
                os.makedirs(output_dir, exist_ok=True)
            self.prs.save(output_path)
            logger.info(f"Presentation saved to {output_path} with a total of {len(self.prs.slides)} slides.")
        except Exception as e:
            logger.error(f"Failed to save presentation to {output_path}: {e}", exc_info=True)
            raise

    # --------------------------------------------------------------------------
    # Internal Helper Methods
    # --------------------------------------------------------------------------

    def _get_slide_layout(self, layout_index: int):
        """Safely retrieves a slide layout by index."""
        try:
            if layout_index < len(self.prs.slide_layouts):
                return self.prs.slide_layouts[layout_index]
            else:
                 logger.warning(f"Slide layout index {layout_index} out of bounds ({len(self.prs.slide_layouts)} layouts available). Using layout 0.")
                 return self.prs.slide_layouts[0]
        except Exception as e:
            logger.error(f"Error getting slide layout {layout_index}: {e}. Using layout 0.", exc_info=True)
            return self.prs.slide_layouts[0]

    def _get_placeholder(self, slide, placeholder_idx: int, placeholder_name: Optional[str] = None):
        """Safely retrieves a placeholder by index or name."""
        try:
            if placeholder_name:
                for shape in slide.placeholders:
                    if shape.name == placeholder_name:
                        return shape
            # Fallback to index if name not found or not provided
            if placeholder_idx < len(slide.placeholders):
                return slide.placeholders[placeholder_idx]
            logger.warning(f"Placeholder index {placeholder_idx} (name: {placeholder_name}) not found on slide.")
            return None
        except Exception as e:
            logger.error(f"Error getting placeholder {placeholder_idx} (name: {placeholder_name}): {e}", exc_info=True)
            return None

    def _apply_text_frame_style(self, text_frame: TextFrame,
                                font_name: Optional[str] = None,
                                font_size: Optional[Pt] = None,
                                font_color: Optional[Tuple[int, int, int]] = None,
                                bold: Optional[bool] = None,
                                italic: Optional[bool] = None,
                                alignment: Optional[PP_ALIGN] = None,
                                vertical_anchor: Optional[MSO_ANCHOR] = MSO_ANCHOR.TOP,
                                word_wrap: Optional[bool] = True,
                                auto_size: Optional[MSO_AUTO_SIZE] = None,
                                line_spacing: Optional[float] = None,
                                space_before: Optional[Pt] = None,
                                space_after: Optional[Pt] = None,
                                apply_to_all_paragraphs: bool = True) -> None:
        """Applies detailed styling to a TextFrame and its paragraphs."""
        if vertical_anchor is not None:
            text_frame.vertical_anchor = vertical_anchor
        if word_wrap is not None:
            text_frame.word_wrap = word_wrap
        if auto_size is not None:
            text_frame.auto_size = auto_size

        target_paragraphs = text_frame.paragraphs if apply_to_all_paragraphs else []
        if not apply_to_all_paragraphs and text_frame.paragraphs:
            target_paragraphs = [text_frame.paragraphs[0]]
        elif not text_frame.paragraphs: # Ensure at least one paragraph exists
             text_frame.add_paragraph()
             target_paragraphs = text_frame.paragraphs

        for p in target_paragraphs:
            if alignment is not None: p.alignment = alignment
            if line_spacing is not None: p.line_spacing = line_spacing
            if space_before is not None: p.space_before = space_before
            if space_after is not None: p.space_after = space_after

            # Apply font styles to the paragraph's default font initially
            # This sets the base for new runs. Then iterate runs if needed.
            font = p.font
            if font_name: font.name = font_name
            if font_size: font.size = font_size
            if font_color: font.color.rgb = RGBColor(*font_color)
            if bold is not None: font.bold = bold
            if italic is not None: font.italic = italic

            # Optionally apply to existing runs as well
            # for run in p.runs:
            #     if font_name: run.font.name = font_name
            #     # ... apply other styles similarly to runs ...

    def _add_styled_paragraph(self, text_frame: TextFrame, text: str,
                              font_name: Optional[str] = None,
                              size: Optional[Pt] = None,
                              color: Optional[Tuple[int, int, int]] = None,
                              bold: Optional[bool] = None,
                              italic: Optional[bool] = None,
                              align: Optional[PP_ALIGN] = None,
                              space_after: Optional[Pt] = None,
                              level: int = 0) -> _Paragraph:
        """Adds a new paragraph with specified styling."""
        p = text_frame.add_paragraph()
        p.text = text
        p.level = level
        font = p.font
        if font_name: font.name = font_name
        if size: font.size = size
        if color: font.color.rgb = RGBColor(*color)
        if bold is not None: font.bold = bold
        if italic is not None: font.italic = italic
        if align: p.alignment = align
        if space_after: p.space_after = space_after
        return p

    def _add_bullet_points(self, text_frame: TextFrame, items: List[str],
                           font_name: Optional[str] = None,
                           size: Optional[Pt] = None,
                           color: Optional[Tuple[int, int, int]] = None,
                           level: int = 0,
                           space_after: Optional[Pt] = Pt(4)) -> None:
        """Adds a list of bullet points to a text frame."""
        for item in items:
            p = self._add_styled_paragraph(
                text_frame, text=item, font_name=font_name, size=size,
                color=color, level=level, space_after=space_after
            )
            # Ensure bullet point character is standard (can be customized further)
            p.font.bullet_char = '•'

    def _add_table(self, slide, left: Inches, top: Inches, width: Inches, height: Inches,
                   data: List[List[Any]],
                   col_widths: Optional[List[Inches]] = None,
                   first_row_header: bool = True,
                   first_col_bold: bool = True,
                   row_striping: bool = True) -> Optional[Shape]:
        """Adds and styles a table."""
        if not data or not data[0]:
            logger.warning("Cannot add table: Data is empty.")
            return None

        rows, cols = len(data), len(data[0])
        try:
            shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            table = shape.table

            # Set column widths
            if col_widths and len(col_widths) == cols:
                for i, cw in enumerate(col_widths):
                    table.columns[i].width = cw
            else: # Default: Even distribution
                default_col_width = Emu(width / cols)
                for i in range(cols):
                    table.columns[i].width = default_col_width

            # Populate and style cells
            for r_idx, row_data in enumerate(data):
                for c_idx, cell_data in enumerate(row_data):
                    cell = table.cell(r_idx, c_idx)
                    cell.text = str(cell_data) # Ensure string conversion
                    tf = cell.text_frame
                    tf.margin_left = Inches(0.08)
                    tf.margin_right = Inches(0.08)
                    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                    tf.word_wrap = True # Allow wrapping

                    is_header = first_row_header and r_idx == 0
                    is_alt_row = row_striping and not is_header and r_idx % 2 == 1
                    is_first_col_data = first_col_bold and not is_header and c_idx == 0

                    # Apply styles
                    self._apply_text_frame_style(
                        tf,
                        font_name=self.design.BODY_FONT,
                        font_size=Pt(10), # Smaller font for tables
                        bold=is_header or is_first_col_data,
                        font_color=self.design.TEXT_LIGHT if is_header else self.design.TEXT_DARK,
                        alignment=PP_ALIGN.LEFT # Default left align
                    )

                    # Cell background
                    if is_header:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(*self.design.TABLE_HEADER_BG)
                    elif is_alt_row:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(*self.design.TABLE_STRIPE_BG)
                    else:
                        cell.fill.background() # No fill / inherit slide background

            return shape
        except Exception as e:
            logger.error(f"Error adding table: {e}", exc_info=True)
            return None

    def _add_textbox_with_style(self, slide, left: Inches, top: Inches, width: Inches, height: Inches,
                                text: str = "", **style_args) -> Optional[Shape]:
         """Adds a textbox and applies initial styling."""
         try:
              textbox = slide.shapes.add_textbox(left, top, width, height)
              textbox.text = text
              self._apply_text_frame_style(textbox.text_frame, **style_args)
              return textbox
         except Exception as e:
              logger.error(f"Error adding styled textbox: {e}", exc_info=True)
              return None

    # --------------------------------------------------------------------------
    # Public Slide Creation Methods (Updated)
    # --------------------------------------------------------------------------
    # Add these missing methods to pptx_builder.py

    def add_fund_overview_slide(self, fund_info: Dict[str, Any]) -> None:
        """
        Creates a fund overview slide with fund details in table format and strategy description.
        
        Args:
            fund_info: Dictionary containing fund information
        """
        title = "Fund Overview"
        
        # Extract fund data for table
        fund_data = []
        for key, value in fund_info.items():
            if key != "strategy" and key != "fund_name":  # Handle these separately
                label = " ".join(word.capitalize() for word in key.split('_'))
                fund_data.append([label, str(value)])
        
        # Extract strategy description
        strategy_description = fund_info.get("strategy", "Strategy information not available.")
        
        # Call the original implementation with the extracted data
        slide_layout = self._get_slide_layout(self.design.TITLE_CONTENT_LAYOUT_IDX)
        slide = self.prs.slides.add_slide(slide_layout)
        self.slide_count += 1
        logger.info(f"Added Fund Overview Slide: '{title}'")
        
        # Set title
        title_shape = self._get_placeholder(slide, 0, "Title 1")
        if title_shape:
            title_shape.text = title
            self._apply_text_frame_style(title_shape.text_frame, font_size=Pt(32), bold=True, font_color=self.design.DARK_COLOR)
        
        # Create layout with table on left and strategy on right
        left_margin = Inches(0.5)
        content_top = Inches(1.3)
        table_width = Inches(5.5)
        
        # Add fund info table
        table_shape = self._add_table(slide, left_margin, content_top, table_width, Inches(4.5), fund_data, 
                                    first_row_header=False, first_col_bold=True)
        
        # Add strategy description on right if provided
        if strategy_description and table_shape:
            strategy_left = left_margin + table_width + Inches(0.4)
            strategy_width = self.prs.slide_width - strategy_left - Inches(0.5)
            strategy_box = self._add_textbox_with_style(
                slide, strategy_left, content_top, strategy_width, Inches(4.5),
                font_name=self.design.BODY_FONT, font_size=Pt(12), font_color=self.design.DARK_COLOR
            )
            if strategy_box:
                tf = strategy_box.text_frame
                tf.clear()
                self._add_styled_paragraph(tf, "Strategy Description", bold=True, size=Pt(14), space_after=Pt(6))
                self._add_styled_paragraph(tf, strategy_description, size=Pt(11))

    def add_team_analysis_slide(self, team_data: Dict[str, Any]) -> None:
        """
        Creates a slide showing key team members with their backgrounds.
        
        Args:
            team_data: Dictionary with team information including key_personnel
        """
        title = "Team Analysis"
        team_profiles = team_data.get("key_personnel", [])
        
        if not team_profiles:
            logger.warning("No team profiles provided. Skipping team analysis slide.")
            return
            
        slide_layout = self._get_slide_layout(self.design.TITLE_CONTENT_LAYOUT_IDX)
        slide = self.prs.slides.add_slide(slide_layout)
        self.slide_count += 1
        logger.info(f"Added Team Analysis Slide: '{title}'")
        
        # Set title
        title_shape = self._get_placeholder(slide, 0, "Title 1")
        if title_shape:
            title_shape.text = title
            self._apply_text_frame_style(title_shape.text_frame, font_size=Pt(32), bold=True, font_color=self.design.DARK_COLOR)
        
        # Calculate layout for team profiles (grid layout)
        profiles_per_row = min(3, len(team_profiles))
        rows_needed = math.ceil(len(team_profiles) / profiles_per_row)
        
        margin = Inches(0.5)
        content_top = Inches(1.3)
        available_width = self.prs.slide_width - (2 * margin)
        available_height = self.prs.slide_height - content_top - margin
        
        profile_width = (available_width - (Inches(0.25) * (profiles_per_row - 1))) / profiles_per_row
        profile_height = (available_height - (Inches(0.25) * (rows_needed - 1))) / rows_needed
        
        # Add team profile textboxes
        for i, profile in enumerate(team_profiles):
            row = i // profiles_per_row
            col = i % profiles_per_row
            
            left = margin + (col * (profile_width + Inches(0.25)))
            top = content_top + (row * (profile_height + Inches(0.25)))
            
            profile_box = self._add_textbox_with_style(
                slide, left, top, profile_width, profile_height,
                font_name=self.design.BODY_FONT, font_size=Pt(11), font_color=self.design.DARK_COLOR
            )
            
            if profile_box:
                tf = profile_box.text_frame
                tf.clear()
                self._add_styled_paragraph(tf, profile.get('name', 'N/A'), bold=True, size=Pt(14), space_after=Pt(2))
                self._add_styled_paragraph(tf, profile.get('title', 'N/A'), italic=True, size=Pt(12), space_after=Pt(6), color=self.design.ACCENT_2)
                
                # Format background into bullet points if not already
                background = profile.get('background', '')
                if background:
                    # Check if background is already bullet-formatted
                    if background.startswith('• '):
                        self._add_styled_paragraph(tf, background, size=Pt(10), space_after=Pt(2))
                    else:
                        # Split on newlines and format as bullets
                        for line in background.split('\n'):
                            if line.strip():
                                self._add_styled_paragraph(tf, line.strip(), size=Pt(10), level=1, space_after=Pt(2))

    # These updated methods should replace the ones in the fixed-code artifact
# They properly use the more comprehensive ChartFactory methods you've provided

    def add_portfolio_allocation_slide(self, portfolio_data: Dict[str, float], market_analysis: Optional[Dict] = None) -> None:
        """
        Creates a slide showing portfolio allocation as a pie chart.
        
        Args:
            portfolio_data: Dictionary with asset allocations
            market_analysis: Optional market analysis data
        """
        title = "Portfolio Allocation"
        
        # Convert portfolio data to format needed for chart
        chart_data = [(asset, value * 100) for asset, value in portfolio_data.items()]
        if not chart_data:
            logger.warning("No chart data provided. Skipping portfolio allocation slide.")
            return
            
        slide_layout = self._get_slide_layout(self.design.TITLE_CONTENT_LAYOUT_IDX)
        slide = self.prs.slides.add_slide(slide_layout)
        self.slide_count += 1
        logger.info(f"Added Portfolio Allocation Slide: '{title}'")
        
        # Set title
        title_shape = self._get_placeholder(slide, 0, "Title 1")
        if title_shape:
            title_shape.text = title
            self._apply_text_frame_style(title_shape.text_frame, font_size=Pt(32), bold=True, font_color=self.design.DARK_COLOR)
        
        # Add pie chart using the chart factory
        chart_left = 1.5
        chart_top = 1.8
        chart_width = 6.0
        chart_height = 4.5
        
        # Convert percentage values from decimal to percentage for display
        display_data = [(item[0], item[1]) for item in chart_data]
        
        # Use chart factory to create pie chart
        try:
            self.chart_factory.create_pie_chart(
                slide, 
                display_data, 
                chart_left, 
                chart_top, 
                chart_width, 
                chart_height, 
                title="Asset Allocation"
            )
            
            # Add table with the same data
            table_data = [["Asset", "Allocation"]]
            for asset, pct in chart_data:
                table_data.append([asset, f"{pct:.1f}%"])
            
            table_left = chart_left + chart_width + 0.5
            table_width = self.prs.slide_width / Inches(1) - table_left - 0.5
            
            self._add_table(slide, Inches(table_left), Inches(chart_top), Inches(table_width), Inches(3.5), table_data)
            
        except Exception as e:
            logger.error(f"Error creating portfolio allocation chart: {e}", exc_info=True)
            # Fallback: just display as table if chart creation fails
            table_data = [["Asset", "Allocation"]]
            for asset, pct in chart_data:
                table_data.append([asset, f"{pct:.1f}%"])
            
            self._add_table(slide, Inches(1.5), Inches(2.0), Inches(9.0), Inches(4.0), table_data)

    def add_wallet_security_analysis(self, wallet_analysis: Dict[str, Any]) -> None:
        """
        Creates a slide showing wallet infrastructure overview.
        
        Args:
            wallet_analysis: Dictionary with wallet analysis results
        """
        title = "Wallet Security Analysis"
        self._add_styled_paragraph(tf, str(total_balance), bold=True, size=Pt(24), align=PP_ALIGN.CENTER, color=self.design.ACCENT_1)
        # Extract data from wallet analysis
        total_balance = wallet_analysis.get("aggregate_stats", {}).get("total_balance_eth", 0)
        total_balance_str = f"{total_balance:.2f} ETH"
        
        wallet_count = len(wallet_analysis.get("wallets", {}))
        
        avg_risk_score = wallet_analysis.get("aggregate_stats", {}).get("average_risk_score", 50)
        
        # Create chart data from wallets
        wallet_types = {}
        for address, wallet in wallet_analysis.get("wallets", {}).items():
            wtype = wallet.get("wallet_type", "Unknown")
            if wtype in wallet_types:
                wallet_types[wtype] += wallet.get("balance", 0)
            else:
                wallet_types[wtype] = wallet.get("balance", 0)
        
        wallet_chart_data = [(wtype, balance) for wtype, balance in wallet_types.items()]
        slide_layout = self._get_slide_layout(self.design.TITLE_CONTENT_LAYOUT_IDX)
        slide = self.prs.slides.add_slide(slide_layout)
        self.slide_count += 1
        logger.info(f"Added Wallet Security Slide: '{title}'")
        
        # Set title
        title_shape = self._get_placeholder(slide, 0, "Title 1")
        if title_shape:
            title_shape.text = title
            self._apply_text_frame_style(title_shape.text_frame, font_size=Pt(32), bold=True, font_color=self.design.DARK_COLOR)
        
        # Layout
        margin = 0.5
        content_top = 1.3
        
        # Top section: Key metrics boxes
        metrics_width = 3.5
        metrics_height = 1.2
        metrics_gap = 0.25
        
        # Wallet Count box
        count_box = self._add_textbox_with_style(
            slide, Inches(margin), Inches(content_top), Inches(metrics_width), Inches(metrics_height),
            font_name=self.design.BODY_FONT, font_size=Pt(12), font_color=self.design.DARK_COLOR,
            vertical_anchor=MSO_ANCHOR.MIDDLE
        )
        if count_box:
            tf = count_box.text_frame
            tf.clear()
            self._add_styled_paragraph(tf, "Number of Wallets", bold=True, size=Pt(12), align=PP_ALIGN.CENTER)
            self._add_styled_paragraph(tf, str(wallet_count), bold=True, size=Pt(24), align=PP_ALIGN.CENTER, color=self.design.ACCENT_1)
        
        # Total Balance box
        balance_left = margin + metrics_width + metrics_gap
        balance_box = self._add_textbox_with_style(
            slide, Inches(balance_left), Inches(content_top), Inches(metrics_width), Inches(metrics_height),
            font_name=self.design.BODY_FONT, font_size=Pt(12), font_color=self.design.DARK_COLOR,
            vertical_anchor=MSO_ANCHOR.MIDDLE
        )
        if balance_box:
            tf = balance_box.text_frame
            tf.clear()
            self._add_styled_paragraph(tf, "Total Balance", bold=True, size=Pt(12), align=PP_ALIGN.CENTER)
            self._add_styled_paragraph(tf, total_balance, bold=True, size=Pt(24), align=PP_ALIGN.CENTER, color=self.design.ACCENT_1)
        
        # Risk Score box - use gauge chart
        score_left = balance_left + metrics_width + metrics_gap
        
        # Bottom section: Wallet types chart
        chart_top = content_top + metrics_height + 0.4
        chart_height = 3.5
        
        if wallet_chart_data:
            # For risk score, use a gauge chart
            try:
                self.chart_factory.create_gauge_chart(
                    slide,
                    avg_risk_score,
                    score_left,
                    content_top,
                    metrics_width,
                    metrics_height,
                    title="Security Risk Score",
                    min_value=0,
                    max_value=100,
                    # Invert thresholds since lower is better for risk
                    green_threshold=30,
                    yellow_threshold=70
                )
            except Exception as e:
                logger.error(f"Error creating risk gauge chart: {e}", exc_info=True)
                # Fallback to text box
                risk_box = self._add_textbox_with_style(
                    slide, Inches(score_left), Inches(content_top), Inches(metrics_width), Inches(metrics_height),
                    font_name=self.design.BODY_FONT, font_size=Pt(12), font_color=self.design.DARK_COLOR,
                    vertical_anchor=MSO_ANCHOR.MIDDLE
                )
                if risk_box:
                    tf = risk_box.text_frame
                    tf.clear()
                    self._add_styled_paragraph(tf, "Average Risk Score", bold=True, size=Pt(12), align=PP_ALIGN.CENTER)
                    risk_color = self.design.get_risk_scale_color(avg_risk_score, invert=True)
                    risk_level = "Low" if avg_risk_score < 40 else "Medium" if avg_risk_score < 70 else "High"
                    self._add_styled_paragraph(tf, f"{avg_risk_score:.1f} ({risk_level})", bold=True, size=Pt(20), 
                                        align=PP_ALIGN.CENTER, color=risk_color)
            
            # For wallet types, use a pie chart
            chart_width = 9.0
            chart_left = 1.5
            
            try:
                # Use chart factory to create pie chart for wallet types
                self.chart_factory.create_pie_chart(
                    slide, 
                    wallet_chart_data, 
                    chart_left, 
                    chart_top, 
                    chart_width, 
                    chart_height, 
                    title="Wallet Types Distribution"
                )
            except Exception as e:
                logger.error(f"Error creating wallet chart: {e}", exc_info=True)
                # Fallback: display as table
                table_data = [["Wallet Type", "Balance"]]
                for wtype, balance in wallet_chart_data:
                    table_data.append([wtype, f"{balance:.2f}"])
                
                self._add_table(slide, Inches(margin), Inches(chart_top), Inches(self.prs.slide_width / Inches(1) - (2 * margin)), Inches(chart_height), table_data)

    def add_risk_overview_slide(self, title: str, overall_risk_score: float, risk_level: str, 
                            risk_color: str, radar_labels: List[str], radar_values: List[float],
                            risk_narrative: str) -> None:
        """
        Creates a risk overview slide with radar chart and narrative.
        
        Args:
            title: Slide title
            overall_risk_score: Overall risk score (0-100)
            risk_level: Risk level description
            risk_color: Hex color for risk level
            radar_labels: Labels for radar chart axes
            radar_values: Values for radar chart points
            risk_narrative: LLM-generated risk narrative
        """
        slide_layout = self._get_slide_layout(self.design.TWO_CONTENT_LAYOUT_IDX)
        slide = self.prs.slides.add_slide(slide_layout)
        self.slide_count += 1
        logger.info(f"Added Risk Overview Slide: '{title}'")
        
        # Set title
        title_shape = self._get_placeholder(slide, 0, "Title 1")
        if title_shape:
            title_shape.text = title
            self._apply_text_frame_style(title_shape.text_frame, font_size=Pt(32), bold=True, font_color=self.design.DARK_COLOR)
        
        # Layout
        left_col_left = 0.5
        left_col_width = 5.5
        right_col_left = left_col_left + left_col_width + 0.5
        right_col_width = self.prs.slide_width / Inches(1) - right_col_left - 0.5
        content_top = 1.3
        
        # Use gauge chart for the risk score
        gauge_top = content_top
        gauge_height = 2.0
        gauge_width = 5.0
        gauge_left = left_col_left + (left_col_width - gauge_width) / 2  # Center the gauge
        
        try:
            # For risk score, invert the green/yellow thresholds (lower is better)
            self.chart_factory.create_gauge_chart(
                slide,
                overall_risk_score,
                gauge_left,
                gauge_top,
                gauge_width,
                gauge_height,
                title="Overall Risk Assessment",
                min_value=0,
                max_value=100,
                # Invert thresholds since lower is better for risk
                green_threshold=30,
                yellow_threshold=70
            )
        except Exception as e:
            logger.error(f"Error creating risk gauge chart: {e}", exc_info=True)
            # Fallback to textbox
            # First, try to convert hex color to RGB for display
            try:
                r, g, b = int(risk_color[0:2], 16), int(risk_color[2:4], 16), int(risk_color[4:6], 16)
                display_color = (r, g, b)
            except ValueError:
                display_color = self.design.DARK_COLOR
                
            score_box = self._add_textbox_with_style(
                slide, Inches(gauge_left), Inches(gauge_top), Inches(gauge_width), Inches(gauge_height),
                font_name=self.design.BODY_FONT, font_size=Pt(14), font_color=self.design.DARK_COLOR,
                vertical_anchor=MSO_ANCHOR.MIDDLE
            )
            
            if score_box:
                tf = score_box.text_frame
                tf.clear()
                self._add_styled_paragraph(tf, "Overall Risk Assessment", bold=True, size=Pt(16), align=PP_ALIGN.CENTER)
                self._add_styled_paragraph(tf, f"{risk_level} ({overall_risk_score:.1f}/100)", 
                                        bold=True, size=Pt(28), align=PP_ALIGN.CENTER,
                                        color=display_color)
        
        # Add radar chart below gauge
        radar_top = gauge_top + gauge_height + 0.3
        radar_height = self.prs.slide_height / Inches(1) - radar_top - 0.5
        
        # Add radar chart for risk components
        if radar_labels and radar_values and len(radar_labels) == len(radar_values):
            try:
                self.chart_factory.create_radar_chart(
                    slide, 
                    radar_labels, 
                    radar_values, 
                    left_col_left, 
                    radar_top, 
                    left_col_width, 
                    radar_height, 
                    title="Risk Components"
                )
            except Exception as e:
                logger.error(f"Error creating risk radar chart: {e}", exc_info=True)
                # Fallback to table
                table_data = [["Risk Component", "Score"]]
                for label, value in zip(radar_labels, radar_values):
                    table_data.append([label, f"{value:.1f}"])
                
                self._add_table(slide, Inches(left_col_left), Inches(radar_top), Inches(left_col_width), Inches(radar_height), table_data)
        
        # Right side: Risk narrative
        narrative_box = self._add_textbox_with_style(
            slide, Inches(right_col_left), Inches(content_top), Inches(right_col_width), 
            Inches(self.prs.slide_height / Inches(1) - content_top - 0.5),
            font_name=self.design.BODY_FONT, font_size=Pt(12), font_color=self.design.DARK_COLOR
        )
        
        if narrative_box:
            tf = narrative_box.text_frame
            tf.clear()
            self._add_styled_paragraph(tf, "Risk Assessment", bold=True, size=Pt(16), space_after=Pt(6))
            self._add_styled_paragraph(tf, risk_narrative, size=Pt(11))
        

    def add_compliance_analysis_slides(self, overview_data: Dict, reg_status_data: Optional[Dict] = None,
                                    gaps_data: Optional[Dict] = None) -> None:
        """
        Creates slides for compliance analysis (overview, regulatory status, and gaps).
        Calls specific compliance slide builders for each component.
        
        Args:
            overview_data: Data for the compliance overview slide
            reg_status_data: Data for the regulatory status slide (optional)
            gaps_data: Data for the compliance gaps slide (optional)
        """
        # Add compliance overview slide
        self.add_compliance_overview_slide(**overview_data)
        
        # Add regulatory status slide if data provided
        if reg_status_data and reg_status_data.get('table_data'):
            self.add_table_slide(**reg_status_data)
        
        # Add compliance gaps slide if data provided
        if gaps_data:
            self.add_text_slide(**gaps_data)

    def add_compliance_overview_slide(self, title: str, overall_score: float, compliance_level: str,
                                    jurisdictions: str, kyc_aml_coverage: float,
                                    compliance_narrative: str) -> None:
        """
        Creates a compliance overview slide with compliance score and narrative.
        
        Args:
            title: Slide title
            overall_score: Overall compliance score (0-100)
            compliance_level: Compliance level description
            jurisdictions: List of jurisdictions as comma-separated string
            kyc_aml_coverage: KYC/AML coverage score (0-100)
            compliance_narrative: LLM-generated compliance narrative
        """
        slide_layout = self._get_slide_layout(self.design.TWO_CONTENT_LAYOUT_IDX)
        slide = self.prs.slides.add_slide(slide_layout)
        self.slide_count += 1
        logger.info(f"Added Compliance Overview Slide: '{title}'")
        
        # Set title
        title_shape = self._get_placeholder(slide, 0, "Title 1")
        if title_shape:
            title_shape.text = title
            self._apply_text_frame_style(title_shape.text_frame, font_size=Pt(32), bold=True, font_color=self.design.DARK_COLOR)
        
        # Layout
        left_col_left = Inches(0.5)
        left_col_width = Inches(5.5)
        right_col_left = left_col_left + left_col_width + Inches(0.5)
        right_col_width = self.prs.slide_width - right_col_left - Inches(0.5)
        content_top = Inches(1.3)
        
        # Left side: Compliance metrics
        # Compliance score box
        score_height = Inches(1.5)
        score_box = self._add_textbox_with_style(
            slide, left_col_left, content_top, left_col_width, score_height,
            font_name=self.design.BODY_FONT, font_size=Pt(14), font_color=self.design.DARK_COLOR,
            vertical_anchor=MSO_ANCHOR.MIDDLE
        )
        
        if score_box:
            tf = score_box.text_frame
            tf.clear()
            self._add_styled_paragraph(tf, "Overall Compliance", bold=True, size=Pt(16), align=PP_ALIGN.CENTER)
            
            # Get color based on compliance score
            compliance_color = self.design.get_risk_scale_color(overall_score, invert=False)
            self._add_styled_paragraph(tf, f"{compliance_level} ({overall_score:.1f}/100)", 
                                    bold=True, size=Pt(28), align=PP_ALIGN.CENTER,
                                    color=compliance_color)
        
        # Add compliance info box
        info_top = content_top + score_height + Inches(0.3)
        info_height = Inches(2.0)
        info_box = self._add_textbox_with_style(
            slide, left_col_left, info_top, left_col_width, info_height,
            font_name=self.design.BODY_FONT, font_size=Pt(12), font_color=self.design.DARK_COLOR
        )
        
        if info_box:
            tf = info_box.text_frame
            tf.clear()
            self._add_styled_paragraph(tf, "Compliance Overview", bold=True, size=Pt(14), space_after=Pt(6))
            
            # Jurisdictions
            self._add_styled_paragraph(tf, "Relevant Jurisdictions:", bold=True, size=Pt(12), space_after=Pt(2))
            self._add_styled_paragraph(tf, jurisdictions, size=Pt(11), space_after=Pt(6))
            
            # KYC/AML coverage with gauge visualization
            self._add_styled_paragraph(tf, "KYC/AML Coverage:", bold=True, size=Pt(12), space_after=Pt(2))
            
            kyc_color = self.design.get_risk_scale_color(kyc_aml_coverage, invert=False)
            self._add_styled_paragraph(tf, f"{kyc_aml_coverage:.1f}% Coverage", 
                                    size=Pt(11), bold=True, color=kyc_color, space_after=Pt(2))
        
        # Right side: Compliance narrative
        narrative_box = self._add_textbox_with_style(
            slide, right_col_left, content_top, right_col_width, 
            self.prs.slide_height - content_top - Inches(0.5),
            font_name=self.design.BODY_FONT, font_size=Pt(12), font_color=self.design.DARK_COLOR
        )
        
        if narrative_box:
            tf = narrative_box.text_frame
            tf.clear()
            self._add_styled_paragraph(tf, "Compliance Assessment", bold=True, size=Pt(16), space_after=Pt(6))
            self._add_styled_paragraph(tf, compliance_narrative, size=Pt(11))

    def add_text_slide(self, title: str, content: str, text_size: int = 12) -> None:
        """
        Creates a simple slide with title and text content.
        
        Args:
            title: Slide title
            content: Text content for the slide
            text_size: Font size for the content
        """
        slide_layout = self._get_slide_layout(self.design.TITLE_CONTENT_LAYOUT_IDX)
        slide = self.prs.slides.add_slide(slide_layout)
        self.slide_count += 1
        logger.info(f"Added Text Slide: '{title}'")
        
        # Set title
        title_shape = self._get_placeholder(slide, 0, "Title 1")
        if title_shape:
            title_shape.text = title
            self._apply_text_frame_style(title_shape.text_frame, font_size=Pt(32), bold=True, font_color=self.design.DARK_COLOR)
        
        # Add content
        content_placeholder = self._get_placeholder(slide, 1, "Content Placeholder 2")
        if content_placeholder:
            content_placeholder.text = content
            self._apply_text_frame_style(
                content_placeholder.text_frame, 
                font_name=self.design.BODY_FONT,
                font_size=Pt(text_size),
                font_color=self.design.DARK_COLOR
            )
        else:
            # Fallback if placeholder not found
            margin = Inches(0.5)
            content_top = Inches(1.3)
            content_box = self._add_textbox_with_style(
                slide, margin, content_top, 
                self.prs.slide_width - (2 * margin),
                self.prs.slide_height - content_top - margin,
                font_name=self.design.BODY_FONT,
                font_size=Pt(text_size),
                font_color=self.design.DARK_COLOR
            )
            if content_box:
                content_box.text = content
    def add_cover_slide(self, title: str, subtitle: str,
                      date: Optional[str] = None,
                      background_color: Optional[Tuple[int, int, int]] = None,
                      logo_path: Optional[str] = None) -> None:
        slide_layout = self._get_slide_layout(self.design.TITLE_LAYOUT_IDX)
        slide = self.prs.slides.add_slide(slide_layout)
        self.slide_count += 1
        logger.info(f"Added Cover Slide: '{title}'")

        bg_color = background_color or self.design.COVER_BACKGROUND
        text_color = self.design.TEXT_LIGHT

        # Set background
        try:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*bg_color)
        except Exception as e:
            logger.error(f"Error setting cover slide background color: {e}", exc_info=True)

        # Set title
        title_shape = self._get_placeholder(slide, 0, "Title 1")
        if title_shape:
            title_shape.text = title
            self._apply_text_frame_style(
                title_shape.text_frame, font_name=self.design.TITLE_FONT, font_size=Pt(44),
                font_color=text_color, bold=True, alignment=PP_ALIGN.CENTER
            )
        else: # Fallback if placeholder missing
             self._add_textbox_with_style(slide, Inches(1), Inches(2), Inches(11.33), Inches(1.5), title,
                                          font_name=self.design.TITLE_FONT, font_size=Pt(44), font_color=text_color,
                                          bold=True, alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

        # Set subtitle
        subtitle_shape = self._get_placeholder(slide, 1, "Subtitle 2") # Common index/name
        if subtitle_shape:
            subtitle_shape.text = subtitle
            self._apply_text_frame_style(
                subtitle_shape.text_frame, font_name=self.design.TITLE_FONT, font_size=Pt(24),
                font_color=text_color, alignment=PP_ALIGN.CENTER
            )
        else: # Fallback
            self._add_textbox_with_style(slide, Inches(1), Inches(3.5), Inches(11.33), Inches(1.0), subtitle,
                                         font_name=self.design.TITLE_FONT, font_size=Pt(24), font_color=text_color,
                                         alignment=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.TOP)

        # Add date
        if date:
             self._add_textbox_with_style(slide, Inches(0.5), self.prs.slide_height - Inches(0.75), Inches(4), Inches(0.5),
                                          date, font_name=self.design.BODY_FONT, font_size=Pt(11),
                                          font_color=text_color, alignment=PP_ALIGN.LEFT)

        # Add logo
        if logo_path and os.path.exists(logo_path):
            logo_width, logo_height = Inches(1.5), Inches(0.75)
            logo_left = self.prs.slide_width - logo_width - Inches(0.5)
            logo_top = self.prs.slide_height - logo_height - Inches(0.25)
            try:
                slide.shapes.add_picture(logo_path, logo_left, logo_top, width=logo_width)
            except Exception as e:
                logger.error(f"Error adding logo '{logo_path}': {e}", exc_info=True)
        elif logo_path:
            logger.warning(f"Logo path specified but not found: {logo_path}")

    def _add_gauge_chart(self, slide, left: Inches, top: Inches, width: Inches, height: Inches,
                         score: float, level: str, color_hex: str, title: str) -> None:
        """
        Adds a simple visual block representing a gauge or score.
        (Note: This is a simplified visual, not a true gauge chart).

        Args:
            slide: The slide object.
            left, top, width, height: Position and dimensions using Inches.
            score: The numerical score (0-100).
            level: Optional qualitative level (e.g., "High", "Adequate").
            color_hex: Background color hex string based on score severity.
            title: Title for the gauge (e.g., "Compliance Score").
        """
        try:
            # Background shape
            rect = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
            )
            rect.fill.solid()
            try:
                rect.fill.fore_color.rgb = RGBColor.from_string(color_hex)
            except ValueError:
                logger.warning(f"Invalid gauge color {color_hex}, using default dark grey.")
                rect.fill.fore_color.rgb = RGBColor(*self.design.DARK_COLOR) # Use tuple for fallback
            rect.line.fill.background() # No border

            # Text Box inside the shape
            text_margin = Inches(0.1)
            text_box = slide.shapes.add_textbox(
                left + text_margin, top + text_margin, width - (2*text_margin), height - (2*text_margin)
            )
            tf = text_box.text_frame
            tf.clear()
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.word_wrap = False # Prevent wrapping for score/level
            # tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT # Can cause issues with alignment

            # Add title
            self._add_styled_paragraph(
                tf, title, font_name=self.design.BODY_FONT, size=Pt(11), bold=True,
                color=self.design.TEXT_LIGHT, align=PP_ALIGN.CENTER, space_after=Pt(2)
            )

            # Add Score
            self._add_styled_paragraph(
                tf, f"{score:.1f}", font_name=self.design.BODY_FONT, size=Pt(18), bold=True,
                color=self.design.TEXT_LIGHT, align=PP_ALIGN.CENTER, space_after=Pt(0)
            )

            # Add Level (optional)
            if level:
                self._add_styled_paragraph(
                    tf, level, font_name=self.design.BODY_FONT, size=Pt(11),
                    color=self.design.TEXT_LIGHT, align=PP_ALIGN.CENTER
                )

        except Exception as e:
            logger.error(f"Error adding gauge chart element for '{title}': {e}", exc_info=True)

    # Rewritten to accept LLM text and structured points separately
    def add_risk_assessment_slides(self, overview_data: Dict, factors_data: Optional[Dict] = None,
                                mitigation_data: Optional[Dict] = None) -> None:
        """
        Creates slides for risk assessment (overview, factors, and mitigations).
        Calls specific risk slide builders for each component.
        
        Args:
            overview_data: Data for the risk overview slide
            factors_data: Data for the risk factors slide (optional)
            mitigation_data: Data for the risk mitigation slide (optional)
        """
        # Add risk overview slide
        self.add_risk_overview_slide(**overview_data)
        
        # Add risk factors slide if data provided
        if factors_data:
            self.add_text_slide(**factors_data)
        
        # Add risk mitigation slide if data provided
        if mitigation_data:
            self.add_text_slide(**mitigation_data)
    def add_executive_summary_slide(self, title: str, fund_name: str, aum: str,
                                 strategy: str,
                                 risk_score: float, risk_level: str, risk_color: Tuple[int, int, int],
                                 compliance_score: float, compliance_level: str, compliance_color: Tuple[int, int, int],
                                 summary_text: str, # LLM generated summary
                                 key_strengths: List[str], # Top points from LLM analysis or fallback
                                 key_concerns: List[str]) -> None:
        slide_layout = self._get_slide_layout(self.design.TWO_CONTENT_LAYOUT_IDX) # Two Content layout
        slide = self.prs.slides.add_slide(slide_layout)
        self.slide_count += 1
        logger.info(f"Added Executive Summary Slide: '{title}'")

        # Set Title
        title_shape = self._get_placeholder(slide, 0, "Title 1")
        if title_shape:
            title_shape.text = title
            self._apply_text_frame_style(title_shape.text_frame, font_size=Pt(32), bold=True, font_color=self.design.DARK_COLOR)

        # --- Left Column ---
        left_col_left = Inches(0.5)
        left_col_width = Inches(6.0)
        top_margin = Inches(1.3)

        # Fund Info Box
        info_height = Inches(1.5)
        info_box = self._add_textbox_with_style(slide, left_col_left, top_margin, left_col_width, info_height,
                                                 font_name=self.design.BODY_FONT, font_size=Pt(12), font_color=self.design.DARK_COLOR,
                                                 vertical_anchor=MSO_ANCHOR.TOP)
        if info_box:
             tf = info_box.text_frame
             tf.clear()
             self._add_styled_paragraph(tf, f"Fund: {fund_name}", bold=True, size=Pt(14))
             self._add_styled_paragraph(tf, f"AUM: {aum}", size=Pt(12))
             self._add_styled_paragraph(tf, f"Strategy: {strategy}", size=Pt(12))

        # Gauges Below Info
        gauge_top = top_margin + info_height + Inches(0.2)
        gauge_height = Inches(1.6)
        gauge_width = Inches(2.8)

        self._add_gauge_chart(slide, left_col_left, gauge_top, gauge_width, gauge_height,
                              risk_score, risk_level, self.design.rgb_to_hex(risk_color), "Risk Assessment")

        comp_gauge_left = left_col_left + gauge_width + Inches(0.2)
        self._add_gauge_chart(slide, comp_gauge_left, gauge_top, gauge_width, gauge_height,
                              compliance_score, compliance_level, self.design.rgb_to_hex(compliance_color), "Compliance")

        # Key Strengths Box (Below Gauges)
        strengths_top = gauge_top + gauge_height + Inches(0.2)
        strengths_height = self.prs.slide_height - strengths_top - Inches(0.5) # Fill rest of height
        strengths_box = self._add_textbox_with_style(slide, left_col_left, strengths_top, left_col_width, strengths_height,
                                                      font_name=self.design.BODY_FONT, font_size=Pt(12), font_color=self.design.DARK_COLOR,
                                                      vertical_anchor=MSO_ANCHOR.TOP)
        if strengths_box:
             tf = strengths_box.text_frame
             tf.clear()
             self._add_styled_paragraph(tf, "Key Strengths", bold=True, size=Pt(14), color=self.design.POSITIVE, space_after=Pt(6))
             self._add_bullet_points(tf, key_strengths, size=Pt(11), color=self.design.DARK_COLOR, space_after=Pt(3))


        # --- Right Column ---
        right_col_left = left_col_left + left_col_width + Inches(0.5)
        right_col_width = self.prs.slide_width - right_col_left - Inches(0.5)

        # LLM Summary Text Box
        summary_height = Inches(3.5) # Allocate good space for summary
        summary_box = self._add_textbox_with_style(slide, right_col_left, top_margin, right_col_width, summary_height,
                                                  font_name=self.design.BODY_FONT, font_size=Pt(12), font_color=self.design.DARK_COLOR,
                                                  vertical_anchor=MSO_ANCHOR.TOP)
        if summary_box:
             tf = summary_box.text_frame
             tf.clear()
             self._add_styled_paragraph(tf, "Summary Assessment", bold=True, size=Pt(14), color=self.design.ACCENT_1, space_after=Pt(6))
             self._add_styled_paragraph(tf, summary_text, size=Pt(12)) # Add LLM text

        # Key Concerns Box (Below Summary)
        concerns_top = top_margin + summary_height + Inches(0.2)
        concerns_height = self.prs.slide_height - concerns_top - Inches(0.5)
        concerns_box = self._add_textbox_with_style(slide, right_col_left, concerns_top, right_col_width, concerns_height,
                                                     font_name=self.design.BODY_FONT, font_size=Pt(12), font_color=self.design.DARK_COLOR,
                                                     vertical_anchor=MSO_ANCHOR.TOP)
        if concerns_box:
            tf = concerns_box.text_frame
            tf.clear()
            self._add_styled_paragraph(tf, "Key Concerns", bold=True, size=Pt(14), color=self.design.NEGATIVE, space_after=Pt(6))
            self._add_bullet_points(tf, key_concerns, size=Pt(11), color=self.design.DARK_COLOR, space_after=Pt(3))


    # (Keep add_fund_overview_slide, add_team_analysis_slide, add_portfolio_allocation_slide, etc.)
    # Modify them slightly if needed to use the new helper methods for consistency.
    # Example modification for add_table_slide:

    def add_table_slide(self, title: str, table_data: List[List[str]],
                      notes: Optional[str] = None) -> None:
        slide_layout = self._get_slide_layout(self.design.TITLE_CONTENT_LAYOUT_IDX)
        slide = self.prs.slides.add_slide(slide_layout)
        self.slide_count += 1
        logger.info(f"Added Table Slide: '{title}'")

        # Set title
        title_shape = self._get_placeholder(slide, 0, "Title 1")
        if title_shape:
            title_shape.text = title
            self._apply_text_frame_style(title_shape.text_frame, font_size=Pt(32), bold=True, font_color=self.design.DARK_COLOR)

        # Table position and size
        table_top = Inches(1.3)
        table_left = Inches(0.5)
        table_width = self.prs.slide_width - Inches(1.0)
        max_table_height = self.prs.slide_height - table_top - Inches(1.0) # Leave room for notes/footer
        notes_height = Inches(0.75)

        # Add table using helper
        table_shape = self._add_table(slide, table_left, table_top, table_width, max_table_height, table_data)

        # Add notes below table if provided
        if notes and table_shape:
             notes_top_pos = table_shape.top + table_shape.height + Inches(0.15)
             # Ensure notes don't go off slide
             actual_notes_height = min(notes_height, self.prs.slide_height - notes_top_pos - Inches(0.25))
             if actual_notes_height > Inches(0.2):
                 notes_box = self._add_textbox_with_style(slide, table_left, notes_top_pos, table_width, actual_notes_height,
                                                          text=f"Notes: {notes}", font_name=self.design.BODY_FONT,
                                                          font_size=Pt(9), italic=True, font_color=self.design.TEXT_MUTED)
    def add_executive_summary_slide(self, title: str, fund_name: str, aum: str,
                                strategy: str,
                                risk_score: float, risk_level: str, risk_color: Tuple[int, int, int],
                                compliance_score: float, compliance_level: str, compliance_color: Tuple[int, int, int],
                                summary_text: str, # LLM generated summary
                                key_strengths: List[str], # Top points from LLM analysis or fallback
                                key_concerns: List[str]) -> None:
        """
        Creates an executive summary slide with fund overview, risk/compliance scores, and key points.
        
        Args:
            title: Slide title
            fund_name: Name of the fund
            aum: Assets under management (formatted)
            strategy: Brief strategy description
            risk_score: Overall risk score (0-100)
            risk_level: Risk level as text
            risk_color: RGB tuple for risk color
            compliance_score: Overall compliance score (0-100)
            compliance_level: Compliance level as text
            compliance_color: RGB tuple for compliance color
            summary_text: LLM-generated executive summary
            key_strengths: List of key strengths
            key_concerns: List of key concerns
        """
        slide_layout = self._get_slide_layout(self.design.TWO_CONTENT_LAYOUT_IDX) # Two Content layout
        slide = self.prs.slides.add_slide(slide_layout)
        self.slide_count += 1
        logger.info(f"Added Executive Summary Slide: '{title}'")

        # Set Title
        title_shape = self._get_placeholder(slide, 0, "Title 1")
        if title_shape:
            title_shape.text = title
            self._apply_text_frame_style(title_shape.text_frame, font_size=Pt(32), bold=True, font_color=self.design.DARK_COLOR)

        # --- Left Column ---
        left_col_left = Inches(0.5)
        left_col_width = Inches(6.0)
        top_margin = Inches(1.3)

        # Fund Info Box
        info_height = Inches(1.5)
        info_box = self._add_textbox_with_style(slide, left_col_left, top_margin, left_col_width, info_height,
                                                font_name=self.design.BODY_FONT, font_size=Pt(12), font_color=self.design.DARK_COLOR,
                                                vertical_anchor=MSO_ANCHOR.TOP)
        if info_box:
            tf = info_box.text_frame
            tf.clear()
            self._add_styled_paragraph(tf, f"Fund: {fund_name}", bold=True, size=Pt(14))
            self._add_styled_paragraph(tf, f"AUM: {aum}", size=Pt(12))
            self._add_styled_paragraph(tf, f"Strategy: {strategy}", size=Pt(12))

        # Gauges Below Info
        gauge_top = top_margin + info_height + Inches(0.2)
        gauge_height = Inches(1.6)
        gauge_width = Inches(2.8)

        # Try using gauge charts via ChartFactory if available
        try:
            # Risk gauge - lower is better
            self.chart_factory.create_gauge_chart(
                slide,
                risk_score,
                left_col_left / Inches(1),
                gauge_top / Inches(1),
                gauge_width / Inches(1),
                gauge_height / Inches(1),
                title="Risk Assessment",
                green_threshold=30,  # Lower is better for risk
                yellow_threshold=70
            )
            
            # Compliance gauge - higher is better
            comp_gauge_left = (left_col_left + gauge_width + Inches(0.2)) / Inches(1)
            self.chart_factory.create_gauge_chart(
                slide,
                compliance_score,
                comp_gauge_left,
                gauge_top / Inches(1),
                gauge_width / Inches(1),
                gauge_height / Inches(1),
                title="Compliance",
                green_threshold=70,  # Higher is better for compliance
                yellow_threshold=30
            )
        except Exception as e:
            # Fallback to simple colored textboxes
            logger.warning(f"Falling back to simple gauges due to: {e}")
            
            # Create risk textbox with color
            risk_box = self._add_textbox_with_style(
                slide, left_col_left, gauge_top, gauge_width, gauge_height,
                font_name=self.design.BODY_FONT, font_size=Pt(12), font_color=self.design.TEXT_LIGHT,
                vertical_anchor=MSO_ANCHOR.MIDDLE
            )
            
            if risk_box:
                # Set background color
                fill = risk_box.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(*risk_color)
                
                # Add text
                tf = risk_box.text_frame
                tf.clear()
                self._add_styled_paragraph(tf, "Risk Assessment", bold=True, size=Pt(12), 
                                        align=PP_ALIGN.CENTER, color=self.design.TEXT_LIGHT)
                self._add_styled_paragraph(tf, f"{risk_score:.1f}", bold=True, size=Pt(20), 
                                        align=PP_ALIGN.CENTER, color=self.design.TEXT_LIGHT)
                self._add_styled_paragraph(tf, risk_level, size=Pt(12), 
                                        align=PP_ALIGN.CENTER, color=self.design.TEXT_LIGHT)
            
            # Create compliance textbox with color
            comp_gauge_left = left_col_left + gauge_width + Inches(0.2)
            comp_box = self._add_textbox_with_style(
                slide, comp_gauge_left, gauge_top, gauge_width, gauge_height,
                font_name=self.design.BODY_FONT, font_size=Pt(12), font_color=self.design.TEXT_LIGHT,
                vertical_anchor=MSO_ANCHOR.MIDDLE
            )
            
            if comp_box:
                # Set background color
                fill = comp_box.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(*compliance_color)
                
                # Add text
                tf = comp_box.text_frame
                tf.clear()
                self._add_styled_paragraph(tf, "Compliance", bold=True, size=Pt(12), 
                                        align=PP_ALIGN.CENTER, color=self.design.TEXT_LIGHT)
                self._add_styled_paragraph(tf, f"{compliance_score:.1f}", bold=True, size=Pt(20), 
                                        align=PP_ALIGN.CENTER, color=self.design.TEXT_LIGHT)
                self._add_styled_paragraph(tf, compliance_level, size=Pt(12), 
                                        align=PP_ALIGN.CENTER, color=self.design.TEXT_LIGHT)

        # Key Strengths Box (Below Gauges)
        strengths_top = gauge_top + gauge_height + Inches(0.2)
        strengths_height = self.prs.slide_height - strengths_top - Inches(0.5) # Fill rest of height
        strengths_box = self._add_textbox_with_style(slide, left_col_left, strengths_top, left_col_width, strengths_height,
                                                    font_name=self.design.BODY_FONT, font_size=Pt(12), font_color=self.design.DARK_COLOR,
                                                    vertical_anchor=MSO_ANCHOR.TOP)
        if strengths_box:
            tf = strengths_box.text_frame
            tf.clear()
            self._add_styled_paragraph(tf, "Key Strengths", bold=True, size=Pt(14), color=self.design.POSITIVE, space_after=Pt(6))
            self._add_bullet_points(tf, key_strengths, size=Pt(11), color=self.design.DARK_COLOR, space_after=Pt(3))


        # --- Right Column ---
        right_col_left = left_col_left + left_col_width + Inches(0.5)
        right_col_width = self.prs.slide_width - right_col_left - Inches(0.5)

        # LLM Summary Text Box
        summary_height = Inches(3.5) # Allocate good space for summary
        summary_box = self._add_textbox_with_style(slide, right_col_left, top_margin, right_col_width, summary_height,
                                                font_name=self.design.BODY_FONT, font_size=Pt(12), font_color=self.design.DARK_COLOR,
                                                vertical_anchor=MSO_ANCHOR.TOP)
        if summary_box:
            tf = summary_box.text_frame
            tf.clear()
            self._add_styled_paragraph(tf, "Summary Assessment", bold=True, size=Pt(14), color=self.design.ACCENT_1, space_after=Pt(6))
            self._add_styled_paragraph(tf, summary_text, size=Pt(12)) # Add LLM text

        # Key Concerns Box (Below Summary)
        concerns_top = top_margin + summary_height + Inches(0.2)
        concerns_height = self.prs.slide_height - concerns_top - Inches(0.5)
        concerns_box = self._add_textbox_with_style(slide, right_col_left, concerns_top, right_col_width, concerns_height,
                                                    font_name=self.design.BODY_FONT, font_size=Pt(12), font_color=self.design.DARK_COLOR,
                                                    vertical_anchor=MSO_ANCHOR.TOP)
        if concerns_box:
            tf = concerns_box.text_frame
            tf.clear()
            self._add_styled_paragraph(tf, "Key Concerns", bold=True, size=Pt(14), color=self.design.NEGATIVE, space_after=Pt(6))
            self._add_bullet_points(tf, key_concerns, size=Pt(11), color=self.design.DARK_COLOR, space_after=Pt(3))
    # --- Add other slide methods similarly, using helpers ---
    # --- Remember to update add_risk_overview_slide, add_compliance_overview_slide, ---
    # --- add_conclusion_slide to accept and display the LLM narratives ---

    # Example: Update add_conclusion_slide
    def add_conclusion_slide(self, title: str, fund_name: str, risk_level: str,
                        risk_score: float, compliance_level: str, compliance_score: float,
                        conclusion_summary: str, # LLM Generated Text
                        strengths: List[str], # Top points
                        concerns: List[str]) -> None: # Top points
        """
        Creates a conclusion slide summarizing the key findings of the report.
        
        Args:
            title: Slide title
            fund_name: Name of the fund
            risk_level: Risk level as text
            risk_score: Overall risk score (0-100)
            compliance_level: Compliance level as text
            compliance_score: Overall compliance score (0-100)
            conclusion_summary: LLM-generated conclusion text
            strengths: List of key strengths (max 3)
            concerns: List of key concerns (max 3)
        """
        slide_layout = self._get_slide_layout(self.design.TWO_CONTENT_LAYOUT_IDX)
        slide = self.prs.slides.add_slide(slide_layout)
        self.slide_count += 1
        logger.info(f"Added Conclusion Slide: '{title}'")

        # Title
        title_shape = self._get_placeholder(slide, 0, "Title 1")
        if title_shape:
            title_shape.text = title
            self._apply_text_frame_style(title_shape.text_frame, font_size=Pt(32), bold=True, font_color=self.design.DARK_COLOR)

        # Layout
        col_width = Inches(6.0)
        col_gap = Inches(0.5)
        col1_left = Inches(0.5)
        col2_left = col1_left + col_width + col_gap
        content_top = Inches(1.3)
        content_height = self.prs.slide_height - content_top - Inches(0.5)

        # Left Column: Summary Text + Scores
        left_box = self._add_textbox_with_style(slide, col1_left, content_top, col_width, content_height,
                                                font_name=self.design.BODY_FONT, font_size=Pt(12), font_color=self.design.DARK_COLOR)
        if left_box:
            tf = left_box.text_frame
            tf.clear()
            self._add_styled_paragraph(tf, f"Overall Assessment: {fund_name}", bold=True, size=Pt(16), space_after=Pt(6))
            self._add_styled_paragraph(tf, conclusion_summary, size=Pt(12), space_after=Pt(12)) # LLM Text

            # Add scores visually
            risk_color = self.design.get_risk_scale_color(risk_score, invert=True)
            compliance_color = self.design.get_risk_scale_color(compliance_score, invert=False)
            
            p_scores = self._add_styled_paragraph(tf, "", size=Pt(12), space_after=Pt(4)) # Empty paragraph for scores
            
            # Add risk score with color
            run = p_scores.add_run()
            run.text = f"Risk: {risk_level} ({risk_score:.1f})"
            run.font.color.rgb = RGBColor(*risk_color)
            run.font.bold = True
            
            # Add separator
            run = p_scores.add_run()
            run.text = "  |  "
            run.font.color.rgb = RGBColor(*self.design.TEXT_MUTED)
            
            # Add compliance score with color
            run = p_scores.add_run()
            run.text = f"Compliance: {compliance_level} ({compliance_score:.1f})"
            run.font.color.rgb = RGBColor(*compliance_color)
            run.font.bold = True

        # Right Column: Strengths and Concerns
        right_box = self._add_textbox_with_style(slide, col2_left, content_top, col_width, content_height,
                                                font_name=self.design.BODY_FONT, font_size=Pt(12), font_color=self.design.DARK_COLOR)
        if right_box:
            tf = right_box.text_frame
            tf.clear()
            # Strengths
            self._add_styled_paragraph(tf, "Key Strengths", bold=True, size=Pt(14), color=self.design.POSITIVE, space_after=Pt(4))
            self._add_bullet_points(tf, strengths, size=Pt(11), space_after=Pt(3))

            # Concerns
            self._add_styled_paragraph(tf, "Key Concerns", bold=True, size=Pt(14), color=self.design.NEGATIVE, space_after=Pt(4), space_before=Pt(12))
            self._add_bullet_points(tf, concerns, size=Pt(11), space_after=Pt(3))