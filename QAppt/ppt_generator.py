import json
from numpy import char
import weaviate
from typing import List, Dict, Optional
import logging
import traceback
from datetime import datetime
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt

logger = logging.getLogger(__name__)

from Sample_Data.vector_store.weaviate_client import get_weaviate_client
from .topic_extractor import TopicExtractor

class QAPresentationGenerator:
    def __init__(self):
        self.prs = Presentation()
        self.output_dir = "reports/presentations"
        
        # Initialize Weaviate client with proper error handling
        try:
            # Get Weaviate client properly
            client = get_weaviate_client()
            if not client:
                raise Exception("Failed to get Weaviate client")
            self.weaviate_client = client
            logger.info("Weaviate client initialized successfully")
        except Exception as e:
            logger.error(f"Error initializing Weaviate client: {e}")
            self.weaviate_client = None
            
        # Create output directory
        os.makedirs(self.output_dir, exist_ok=True)
        
        # Define colors
        self.colors = {
            'primary': RGBColor(44, 62, 80),
            'secondary': RGBColor(52, 152, 219),
            'accent1': RGBColor(46, 204, 113),
            'accent2': RGBColor(231, 76, 60),
            'text': RGBColor(44, 62, 80)
        }
        
        # Initialize topic extractor with proper error handling
        try:
            self.topic_extractor = TopicExtractor()
            logger.info("Topic extractor initialized successfully")
        except Exception as e:
            logger.error(f"Error initializing topic extractor: {e}")
            self.topic_extractor = None

    def __del__(self):
        """Cleanup method"""
        try:
            # Close Weaviate client
            if hasattr(self, 'weaviate_client') and self.weaviate_client:
                self.weaviate_client.close()
                
            # Clean up topic extractor
            if hasattr(self, 'topic_extractor'):
                del self.topic_extractor
                
        except Exception as e:
            logger.error(f"Error cleaning up presentation generator: {e}")

    def create_title_slide(self, title: str, subtitle: str = ""):
        """Create the title slide for the presentation"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])  # Using title slide layout
        
        # Add and format title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_frame = title_shape.text_frame
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        
        # Add and format subtitle
        subtitle_shape = slide.placeholders[1]  # Index 1 is subtitle in title slide layout
        subtitle_shape.text = f"{subtitle}\nGenerated on: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
        subtitle_frame = subtitle_shape.text_frame
        subtitle_frame.paragraphs[0].font.size = Pt(24)
        subtitle_frame.paragraphs[0].font.color.rgb = self.colors['secondary']

    def add_section_slide(self, title: str):
        """Add a section divider slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[2])  # Using section header layout
        title_shape = slide.shapes.title
        title_shape.text = title
        title_frame = title_shape.text_frame
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.color.rgb = self.colors['secondary']

    def add_content_slide(self, title: str, content: list):
        """Add a content slide with bullet points"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])  # Using content layout
        
        # Add title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Add content
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        
        for point in content:
            p = tf.add_paragraph()
            p.text = str(point)
            p.font.size = Pt(18)
            p.font.color.rgb = self.colors['text']
            p.level = 0

    def add_visualization_slide(self, title: str, data: dict, chart_type: str = 'line'):
        """Add a slide with visualization"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])  # Blank layout
        
        # Add title
        title_box = slide.shapes.title
        title_box.text = title
        title_box.text_frame.paragraphs[0].font.size = Pt(32)
        title_box.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']

        # Create visualization
        plt.figure(figsize=(10, 6))
        if chart_type == 'line':
            plt.plot(data['x'], data['y'])
        elif chart_type == 'bar':
            plt.bar(data['x'], data['y'])
        elif chart_type == 'pie':
            plt.pie(data['values'], labels=data['labels'])

        plt.title(title)
        
        # Save temp image
        temp_image = os.path.join(self.output_dir, 'temp_chart.png')
        plt.savefig(temp_image)
        plt.close()

        # Add image to slide
        image_left = Inches(2)
        image_top = Inches(2)
        slide.shapes.add_picture(temp_image, image_left, image_top, height=Inches(4))
        
        # Clean up temp file
        os.remove(temp_image)

    def add_key_metrics_slide(self, title: str, metrics: dict):
        """Add a slide with key metrics in a grid layout"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        
        # Add title
        title_box = slide.shapes.title
        title_box.text = title
        title_box.text_frame.paragraphs[0].font.size = Pt(32)
        
        # Create grid of metrics
        left = top = width = height = Inches(2)
        for i, (metric, value) in enumerate(metrics.items()):
            row = i // 2
            col = i % 2
            
            # Add metric box
            left = Inches(1 + col * 4.5)
            top = Inches(2 + row * 1.5)
            
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            
            # Add metric name
            p = tf.add_paragraph()
            p.text = metric
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = self.colors['secondary']
            
            # Add value
            p = tf.add_paragraph()
            p.text = str(value)
            p.font.size = Pt(28)
            p.font.color.rgb = self.colors['primary']

    def add_summary_slide(self, title: str, summary_points: list):
        """Add a slide with key summary points"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        
        # Add title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Add summary points
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        
        for point in summary_points:
            p = tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(18)
            p.font.color.rgb = self.colors['text']

    def generate_report(self, conversation_history: List[Dict]) -> str:
        """Generate a comprehensive report based on conversation history"""
        try:
            # Extract topics from conversation
            topics = self.topic_extractor.extract_topics(conversation_history)
            
            # Create title slide
            self.create_title_slide(
                "Crypto Analysis Report",
                "Comprehensive Overview of Discussed Topics"
            )

            # Add topics overview slide
            self.add_summary_slide(
                "Discussion Topics",
                [
                    f"Main Topics: {', '.join(topics['main_topics'])}",
                    f"Technical Areas: {', '.join(topics['technical_topics'])}",
                    f"Market Analysis: {', '.join(topics['market_topics'])}",
                    f"Related Entities: {', '.join(topics['entities'])}"
                ]
            )

            # Generate slides for each main topic
            for topic in topics['main_topics']:
                self._add_topic_analysis_slides(topic, conversation_history)

            # Add technical analysis if present
            if topics['technical_topics']:
                self._add_technical_analysis_slides(topics['technical_topics'])

            # Add market analysis if present
            if topics['market_topics']:
                self._add_market_analysis_slides(topics['market_topics'])

            # Save the presentation
            output_file = os.path.join(
                self.output_dir,
                f"crypto_analysis_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            )
            self.prs.save(output_file)
            return output_file

        except Exception as e:
            logger.error(f"Error generating report: {e}")
            return None

    def generate_comprehensive_report(self, topics: List[str], topic_discussions: Dict, conversation_history: List[Dict]) -> str:
        """Generate a comprehensive report about all discussed topics"""
        try:
            # Create title slide
            self.create_title_slide("Cryptocurrency Analysis Report", "Based on Q&A Discussion")

            # Add overview slide
            overview_content = [
                f"Total Topics Discussed: {len(topics)}",
                f"Topics: {', '.join(topics)}",
                f"Time Period: {conversation_history[0]['timestamp'].strftime('%Y-%m-%d %H:%M')} to {conversation_history[-1]['timestamp'].strftime('%Y-%m-%d %H:%M')}"
            ]
            self.add_content_slide("Discussion Overview", overview_content)

            # Process each topic
            for topic in topics:
                if topic in topic_discussions:
                    # Add topic section header
                    self.add_section_slide(f"{topic.title()} Analysis")

                    # Get discussions for this topic
                    discussions = topic_discussions[topic]

                    # Add key points from discussions
                    key_points = []
                    for entry in discussions:
                        q = entry['question']
                        a = entry['answer'][:200] + "..." if len(entry['answer']) > 200 else entry['answer']
                        key_points.append(f"Q: {q}")
                        key_points.append(f"A: {a}")

                    self.add_content_slide(f"{topic.title()} Discussion Points", key_points)

                    # Try to add market data if available
                    try:
                        market_data = self._fetch_market_data(topic)
                        if market_data:
                            self.add_visualization_slide(f"{topic} Market Performance", market_data)
                    except Exception as e:
                        logger.error(f"Error adding market data for {topic}: {e}")

            # Save the presentation
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            output_file = os.path.join(self.output_dir, f"crypto_analysis_report_{timestamp}.pptx")
            self.prs.save(output_file)
            return output_file

        except Exception as e:
            logger.error(f"Error generating comprehensive report: {e}")
            return None

    def _fetch_market_data(self, topic: str) -> Dict:
        """Fetch market data from Weaviate"""
        try:
            # Get the collection properly
            collection = self.weaviate_client.collections.get("MarketMetrics")
            
            # Query using proper syntax
            response = collection.query.fetch_objects(
                limit=30,
                where={
                    "path": ["symbol"],
                    "operator": "Like",
                    "valueString": topic.upper()
                }
            )

            if response and response.objects:
                # Format data for visualization
                data_points = []
                for obj in response.objects:
                    data_point = {
                        "timestamp": obj.properties.get("timestamp", ""),
                        "price": float(obj.properties.get("price", 0)),
                        "volume": float(obj.properties.get("volume_24h", 0)),
                        "market_cap": float(obj.properties.get("market_cap", 0)),
                        "change_24h": float(obj.properties.get("price_change_24h", 0))
                    }
                    data_points.append(data_point)

                # Sort by timestamp
                data_points.sort(key=lambda x: x["timestamp"])

                return {
                    "x": [point["timestamp"] for point in data_points],
                    "y": [point["price"] for point in data_points],
                    "type": "line",
                    "title": f"{topic.upper()} Price History"
                }

            return None
        except Exception as e:
            logger.error(f"Error fetching market data: {e}")
            return None

    def _extract_key_facts(self, documents):
        """Extract key facts from documents"""
        facts = []
        for doc in documents.get("data", {}).get("Get", {}).get("CryptoDueDiligenceDocuments", []):
            content = doc.get("content", "")
            # Extract sentences that contain numbers or key statistics
            for sentence in content.split("."):
                if any(char.isdigit() for sentence in sentence):
                    facts.append(sentence.strip())
        return facts[:5]  # Return top 5 facts

    def _summarize_qa_conversation(self, content):
        """Summarize the Q&A conversation"""
        summary = []
        qa_pairs = content.split("Question:")
        for qa in qa_pairs:
            if "Answer:" in qa:
                q, a = qa.split("Answer:")
                summary.append(f"Q: {q.strip()}")
                # Extract key points from answer
                key_points = [p.strip() for p in a.split("\n") if p.strip() and not p.startswith("*")]
                summary.extend([f"• {p}" for p in key_points[:3]])
        return summary

    def _analyze_sentiment(self, sentiment_data):
        """Analyze sentiment distribution"""
        sentiments = []
        for item in sentiment_data.get("data", {}).get("Get", {}).get("CryptoNewsSentiment", []):
            sentiments.append(item.get("sentiment", "neutral"))
        
        # Count sentiment distribution
        positive = sentiments.count("positive")
        negative = sentiments.count("negative")
        neutral = sentiments.count("neutral")
        
        return {
            "values": [positive, neutral, negative],
            "labels": ["Positive", "Neutral", "Negative"]
        }

    def _extract_technical_data(self, documents):
        """Extract technical analysis points"""
        technical_points = []
        keywords = ["protocol", "blockchain", "technology", "implementation", "architecture"]
        
        for doc in documents.get("data", {}).get("Get", {}).get("CryptoDueDiligenceDocuments", []):
            content = doc.get("content", "")
            for sentence in content.split("."):
                if any(keyword in sentence.lower() for keyword in keywords):
                    technical_points.append(sentence.strip())
        
        return technical_points[:5]

    def _extract_risks(self, documents):
        """Extract risk factors"""
        risks = []
        keywords = ["risk", "threat", "vulnerability", "concern", "warning"]
        
        for doc in documents.get("data", {}).get("Get", {}).get("CryptoDueDiligenceDocuments", []):
            content = doc.get("content", "")
            for sentence in content.split("."):
                if any(keyword in sentence.lower() for keyword in keywords):
                    risks.append(sentence.strip())
        
        return risks[:5]
    
    def _add_topic_analysis_slides(self, topic: str, conversation_history: List[Dict]):
        """Add slides analyzing a specific topic with visuals and insights"""
        # Add topic section header
        self.add_section_slide(f"{topic} Analysis")
        
        # Extract topic-specific Q&A from conversation history
        topic_qa = []
        for entry in conversation_history:
            # Check if this topic appears in question or answer
            if (topic.lower() in entry['question'].lower() or 
                topic.lower() in entry['answer'].lower()):
                topic_qa.append(entry)
        
        # Add key Q&A points
        qa_points = []
        for entry in topic_qa[:3]:  # Limit to top 3 for clarity
            q = entry['question']
            a = entry['answer'][:150] + "..." if len(entry['answer']) > 150 else entry['answer']
            qa_points.append(f"Q: {q}")
            qa_points.append(f"A: {a}")
        
        if qa_points:
            self.add_content_slide(f"{topic} Discussion Points", qa_points)
        
        # Fetch market data from Weaviate
        market_data = self._fetch_market_data(topic)
        if market_data:
            self.add_visualization_slide(f"{topic} Market Performance", market_data)
        
        # Fetch key metrics from Weaviate
        metrics = self._fetch_key_metrics(topic)
        if metrics:
            self.add_key_metrics_slide(f"{topic} Key Metrics", metrics)
        
        # Fetch sentiment analysis
        sentiment_data = self._fetch_sentiment_data(topic)
        if sentiment_data:
            self.add_visualization_slide(f"{topic} Sentiment Analysis", 
                                    sentiment_data, chart_type='pie')
        
        # Fetch related documents and extract insights
        insights = self._fetch_topic_insights(topic)
        if insights:
            self.add_content_slide(f"{topic} Key Insights", insights)

    def _add_technical_analysis_slides(self, technical_topics: List[str]):
        """Add slides with technical analysis of cryptocurrencies"""
        # Add section header
        self.add_section_slide("Technical Analysis")
        
        # Create a summary of technical topics
        tech_summary = [f"Topics covered: {', '.join(technical_topics)}"]
        self.add_content_slide("Technical Topics Overview", tech_summary)
        
        # For each technical topic, add detailed analysis
        for topic in technical_topics[:3]:  # Limit to top 3 for clarity
            technical_data = self._fetch_technical_data(topic)
            if technical_data:
                self.add_content_slide(f"{topic} Technical Details", technical_data)
        
        # Add architecture diagram if available
        architecture_data = self._fetch_architecture_data(technical_topics[0])
        if architecture_data:
            self.add_visualization_slide("Architecture Overview", architecture_data)

    def _add_market_analysis_slides(self, market_topics: List[str]):
        """Add slides with market analysis"""
        # Add section header
        self.add_section_slide("Market Analysis")
        
        # Market overview slide
        market_overview = self._fetch_market_overview()
        if market_overview:
            self.add_content_slide("Market Overview", market_overview)
        
        # Add price comparison chart
        comparison_data = self._fetch_price_comparison(market_topics)
        if comparison_data:
            self.add_visualization_slide("Price Comparison", comparison_data)
        
        # Add trading volume analysis
        volume_data = self._fetch_trading_volume(market_topics)
        if volume_data:
            self.add_visualization_slide("Trading Volume", volume_data, chart_type='bar')
        
        # Add market correlations
        correlation_data = self._fetch_correlations(market_topics)
        if correlation_data:
            self.add_content_slide("Market Correlations", correlation_data)

    def _fetch_key_metrics(self, topic: str) -> Dict:
        """Fetch key metrics for a topic from Weaviate"""
        try:
            # Get the collection
            collection = self.weaviate_client.collections.get("TokenMetrics")
            
            # Query using proper syntax
            response = collection.query.fetch_objects(
                limit=1,
                where={
                    "path": ["symbol"],
                    "operator": "Like",
                    "valueString": topic.upper()
                }
            )

            if response and response.objects:
                # Get the latest metrics
                metrics = response.objects[0].properties
                
                # Format metrics for display
                return {
                    "Market Cap": f"${float(metrics.get('market_cap', 0)):,.2f}",
                    "24h Volume": f"${float(metrics.get('volume_24h', 0)):,.2f}",
                    "Circulating Supply": f"{float(metrics.get('circulating_supply', 0)):,.0f}",
                    "Max Supply": f"{float(metrics.get('max_supply', 0)):,.0f}",
                    "24h Change": f"{float(metrics.get('price_change_24h', 0)):.2f}%"
                }
            return None
        except Exception as e:
            logger.error(f"Error fetching key metrics: {e}")
            return None

    def _fetch_sentiment_data(self, topic: str) -> Dict:
        """Fetch sentiment analysis data from Weaviate"""
        try:
            # Get the collection
            collection = self.weaviate_client.collections.get("CryptoNewsSentiment")
            
            # Query using proper syntax
            response = collection.query.fetch_objects(
                limit=50,
                where={
                    "path": ["related_tokens"],
                    "operator": "ContainsAny",
                    "valueStringArray": [topic.upper()]
                }
            )

            if response and response.objects:
                # Count sentiment distribution
                positive = 0
                negative = 0 
                neutral = 0
                
                for item in response.objects:
                    sentiment = item.properties.get("sentiment", "neutral")
                    if sentiment == "positive":
                        positive += 1
                    elif sentiment == "negative":
                        negative += 1
                    else:
                        neutral += 1
                
                # Format for pie chart
                return {
                    "labels": ["Positive", "Neutral", "Negative"],
                    "values": [positive, neutral, negative]
                }
            return None
        except Exception as e:
            logger.error(f"Error fetching sentiment data: {e}")
            return None

    def _fetch_topic_insights(self, topic: str) -> List[str]:
        """Fetch key insights about a topic from Weaviate"""
        try:
            # Get the collection
            collection = self.weaviate_client.collections.get("CryptoDueDiligenceDocuments")
            
            # Query using proper syntax with vector search
            response = collection.query.hybrid(
                query=f"{topic} key insights analysis",
                limit=5,
                where={
                    "path": ["content"],
                    "operator": "Like",
                    "valueString": f"%{topic}%"
                }
            )

            if response and response.objects:
                # Extract key insights from documents
                insights = []
                for doc in response.objects:
                    content = doc.properties.get("content", "")
                    
                    # Find sentences containing the topic
                    for sentence in content.split("."):
                        if topic.lower() in sentence.lower() and len(sentence) > 20:
                            insights.append(sentence.strip())
                
                # Return top 5 unique insights
                return list(set(insights))[:5]
            return None
        except Exception as e:
            logger.error(f"Error fetching topic insights: {e}")
            return None

    def _fetch_technical_data(self, topic: str) -> List[str]:
        """Fetch technical information about a topic from Weaviate"""
        try:
            # Get the collection 
            collection = self.weaviate_client.collections.get("CryptoTechnicalDocuments")
            
            # Query using proper syntax
            response = collection.query.hybrid(
                query=f"{topic} technical architecture protocol",
                limit=3
            )

            if response and response.objects:
                # Extract technical details
                tech_details = []
                for doc in response.objects:
                    content = doc.properties.get("content", "")
                    
                    # Extract key technical sentences
                    tech_keywords = ["protocol", "blockchain", "architecture", "consensus", "network"]
                    for sentence in content.split("."):
                        if any(keyword in sentence.lower() for keyword in tech_keywords):
                            tech_details.append(sentence.strip())
                
                # Return top 5 unique technical details
                return list(set(tech_details))[:5]
            return None
        except Exception as e:
            logger.error(f"Error fetching technical data: {e}")
            return None

    def _fetch_architecture_data(self, topic: str) -> Dict:
        """Create an architecture diagram for visualization"""
        # This would typically fetch real data, but we'll create placeholder data
        return {
            "x": ["Components", "Network", "Consensus", "Security"],
            "y": [85, 70, 90, 65],
            "type": "bar",
            "title": f"{topic} Architecture Analysis"
        }

    def _fetch_market_overview(self) -> List[str]:
        """Fetch market overview data from Weaviate"""
        try:
            # Get the collection
            collection = self.weaviate_client.collections.get("MarketAnalysisReports")
            
            # Query latest market reports
            response = collection.query.fetch_objects(
                limit=1,
                order_by=[{"path": ["timestamp"], "order": "desc"}]
            )

            if response and response.objects:
                # Extract market overview points
                overview = response.objects[0].properties.get("summary", "").split("\n")
                return [point.strip() for point in overview if point.strip()]
            return ["No market overview data available"]
        except Exception as e:
            logger.error(f"Error fetching market overview: {e}")
            return ["Error fetching market data"]

    def _fetch_price_comparison(self, market_topics: List[str]) -> Dict:
        """Fetch price comparison data for multiple currencies"""
        try:
            prices = {}
            dates = []
            
            # For each topic, get price history
            for topic in market_topics[:3]:  # Limit to top 3
                # Get collection
                collection = self.weaviate_client.collections.get("MarketMetrics")
                
                # Query price data
                response = collection.query.fetch_objects(
                    limit=7,  # Last 7 data points
                    where={
                        "path": ["symbol"],
                        "operator": "Like",
                        "valueString": topic.upper()
                    },
                    order_by=[{"path": ["timestamp"], "order": "desc"}]
                )
                
                if response and response.objects:
                    # Extract price data
                    topic_prices = []
                    
                    for obj in response.objects:
                        price = float(obj.properties.get("price", 0))
                        timestamp = obj.properties.get("timestamp", "")
                        
                        topic_prices.append(price)
                        if timestamp not in dates:
                            dates.append(timestamp)
                    
                    prices[topic.upper()] = topic_prices
            
            # Format for visualization
            if prices and dates:
                return {
                    "x": dates,
                    "y": prices,
                    "type": "line",
                    "title": "Price Comparison"
                }
            return None
        except Exception as e:
            logger.error(f"Error fetching price comparison: {e}")
            return None

    def _fetch_trading_volume(self, market_topics: List[str]) -> Dict:
        """Fetch trading volume data for visualization"""
        try:
            volumes = []
            labels = []
            
            # For each topic, get latest volume
            for topic in market_topics[:5]:  # Limit to top 5
                # Get collection
                collection = self.weaviate_client.collections.get("MarketMetrics")
                
                # Query volume data
                response = collection.query.fetch_objects(
                    limit=1,
                    where={
                        "path": ["symbol"],
                        "operator": "Like",
                        "valueString": topic.upper()
                    },
                    order_by=[{"path": ["timestamp"], "order": "desc"}]
                )
                
                if response and response.objects:
                    # Extract volume data
                    volume = float(response.objects[0].properties.get("volume_24h", 0))
                    volumes.append(volume)
                    labels.append(topic.upper())
            
            # Format for bar chart
            if volumes and labels:
                return {
                    "x": labels,
                    "y": volumes,
                    "type": "bar",
                    "title": "24h Trading Volume"
                }
            return None
        except Exception as e:
            logger.error(f"Error fetching trading volume: {e}")
            return None

    def _fetch_correlations(self, market_topics: List[str]) -> List[str]:
        """Fetch correlation insights between different crypto assets"""
        try:
            # Get the collection
            collection = self.weaviate_client.collections.get("MarketCorrelations")
            
            correlations = []
            
            # For each pair of topics
            for i, topic1 in enumerate(market_topics[:3]):
                for topic2 in market_topics[i+1:4]:  # Create pairs with limit
                    # Query correlations
                    response = collection.query.fetch_objects(
                        limit=1,
                        where={
                            "path": ["pair"],
                            "operator": "Like",
                            "valueString": f"{topic1.upper()}-{topic2.upper()}"
                        }
                    )
                    
                    if response and response.objects:
                        # Extract correlation data
                        corr_value = float(response.objects[0].properties.get("correlation", 0))
                        correlations.append(f"{topic1.upper()} - {topic2.upper()}: {corr_value:.2f} correlation")
            
            # Add some general correlation insights
            if not correlations:
                correlations = [
                    "BTC and ETH typically show strong positive correlation",
                    "DeFi tokens often move together during market trends",
                    "Stablecoins show negative correlation with market volatility",
                    "Layer-1 blockchains often correlate with their ecosystem tokens"
                ]
                
            return correlations
        except Exception as e:
            logger.error(f"Error fetching correlations: {e}")
            return ["Correlation data not available"]