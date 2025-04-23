import re
from pptx import Presentation
from pptx.util import Inches, Pt
import matplotlib.pyplot as plt
from io import BytesIO
from datetime import datetime
import json

# Detect if the user wants to generate a report
def detect_report_intent(question: str) -> bool:
    triggers = [
        r"\bgenerate\b.*\breport\b",
        r"\breport\b.*\bsummary\b",
        r"\bcreate\b.*\breport\b",
        r"\bexport\b.*\bdashboard\b",
        r"\breport\b.*\bnow\b",
        r"\bdashboard\b.*\bplease\b"
    ]
    return any(re.search(pattern, question.lower()) for pattern in triggers)

# Sentiment pie chart
def add_sentiment_slide(prs, system):
    try:
        sentiment = system.get_sentiment_analysis("bitcoin")
        dist = sentiment.get("sentiment_distribution", {})
        if not dist:
            raise ValueError("No sentiment distribution available")
        labels, sizes = list(dist.keys()), list(dist.values())

        fig, ax = plt.subplots()
        ax.pie(sizes, labels=labels, autopct='%1.1f%%', startangle=140)
        ax.axis('equal')

        stream = BytesIO()
        plt.savefig(stream, format='png')
        plt.close(fig)
        stream.seek(0)

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Sentiment Distribution for Bitcoin"
        slide.shapes.add_picture(stream, Inches(2), Inches(1.5), Inches(6), Inches(4.5))
    except Exception:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Sentiment Data"
        tf = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4)).text_frame
        tf.text = "No sentiment data available."

# Risk bar chart
def add_risk_bar_slide(prs, system):
    try:
        docs = system.search("risk", collection="CryptoDueDiligenceDocuments", limit=10)
        entities = [doc.get("title", "Unknown")[:30] for doc in docs]
        scores = [doc.get("risk_score", 0) for doc in docs if doc.get("risk_score") is not None]

        if not entities or not scores:
            raise ValueError("No risk data found")

        fig, ax = plt.subplots()
        ax.barh(entities, scores, color='crimson')
        ax.set_title("Top Risk Scores")
        ax.invert_yaxis()

        stream = BytesIO()
        plt.tight_layout()
        plt.savefig(stream, format='png')
        plt.close(fig)
        stream.seek(0)

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Market Risk Overview"
        slide.shapes.add_picture(stream, Inches(1), Inches(1.5), Inches(8), Inches(4.5))
    except Exception:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Risk Overview"
        tf = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4)).text_frame
        tf.text = "No risk data available."

# Forecast trends
def add_forecast_slide(prs, system):
    for symbol in ["BTCUSDT", "ETHUSDT", "SOLUSDT"]:
        try:
            data = system.get_historical_data(symbol, limit=7)
            if not data:
                raise ValueError("No historical data")
            dates = [d['timestamp'][:10] for d in data]
            prices = [d['close'] for d in data]

            fig, ax = plt.subplots()
            ax.plot(dates, prices, marker='o')
            ax.set_title(f"7-Day Price Trend for {symbol}")
            ax.set_ylabel("Price (USDT)")
            ax.tick_params(axis='x', rotation=45)

            stream = BytesIO()
            plt.tight_layout()
            plt.savefig(stream, format='png')
            plt.close(fig)
            stream.seek(0)

            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = f"Forecast Trends for {symbol}"
            slide.shapes.add_picture(stream, Inches(1), Inches(1.5), Inches(8), Inches(4.5))
        except Exception:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = f"Forecast for {symbol}"
            tf = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4)).text_frame
            tf.text = f"No forecast data available for {symbol}."

# Jurisdictions
def add_jurisdictions_slide(prs, system):
    try:
        docs = system.search("jurisdiction", collection="CryptoDueDiligenceDocuments", limit=20)
        counts = {}
        for doc in docs:
            for word in doc.get("keywords", []):
                key = word.lower()
                if key in ["us", "usa", "united states", "eu", "europe", "uk", "canada", "singapore"]:
                    counts[key] = counts.get(key, 0) + 1

        if not counts:
            raise ValueError("No jurisdiction mentions")

        labels, values = zip(*counts.items())
        fig, ax = plt.subplots()
        ax.bar(labels, values, color='steelblue')
        ax.set_title("Top Mentioned Jurisdictions")
        plt.xticks(rotation=30)

        stream = BytesIO()
        plt.tight_layout()
        plt.savefig(stream, format='png')
        plt.close(fig)
        stream.seek(0)

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Jurisdiction Mentions"
        slide.shapes.add_picture(stream, Inches(1), Inches(1.5), Inches(8), Inches(4.5))
    except Exception:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Jurisdiction Mentions"
        tf = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4)).text_frame
        tf.text = "No jurisdiction data available."

# GPT-style summary from last Q&A
last_answer = ""

def add_summary_slide(prs, question):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "🧠 Key Insights Summary"
    slide.placeholders[1].text = last_answer[:1900] if last_answer else "Summary not available."

# Build the full report
def build_dashboard_report(system, question: str) -> str:
    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Crypto Due Diligence Report"
    slide.placeholders[1].text = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\nQuery: {question}"

    add_sentiment_slide(prs, system)
    add_risk_bar_slide(prs, system)
    add_forecast_slide(prs, system)
    add_jurisdictions_slide(prs, system)
    add_summary_slide(prs, question)

    filename = f"crypto_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    prs.save(filename)
    return filename

# Hook patch
def patched_answer_question(original_answer_fn, system):
    def wrapped(question, *args, **kwargs):
        global last_answer
        if detect_report_intent(question):
            report_path = build_dashboard_report(system, question)
            return f"📊 Report generated and saved to `{report_path}`"

        answer = original_answer_fn(question, *args, **kwargs)
        last_answer = answer
        return answer
    return wrapped